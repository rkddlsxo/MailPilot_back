# services/attachment_service.py - 완전 수정 버전

import os
import io
import tempfile
import hashlib
from pathlib import Path
import numpy as np

# 선택적 임포트 - 없는 라이브러리는 비활성화
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("[⚠️ PIL/Pillow 없음 - 이미지 처리 비활성화]")

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False
    print("[⚠️ pdfplumber 없음 - PDF 처리 비활성화]")

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False
    print("[⚠️ PyPDF2 없음 - PDF 백업 처리 비활성화]")

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("[⚠️ python-docx 없음 - Word 문서 처리 비활성화]")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("[⚠️ python-pptx 없음 - PowerPoint 처리 비활성화]")

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False
    print("[⚠️ pandas 없음 - Excel 처리 비활성화]")

try:
    from pdf2image import convert_from_bytes
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    print("[⚠️ pdf2image 없음 - PDF OCR 처리 비활성화]")

class AttachmentService:
    def __init__(self, config, ai_models):
        self.config = config
        self.ai_models = ai_models
        self.attachment_cache = {}
        
        # 사용 가능한 기능 체크
        self.features = {
            'image_processing': PIL_AVAILABLE,
            'pdf_processing': PDFPLUMBER_AVAILABLE or PYPDF2_AVAILABLE,
            'docx_processing': DOCX_AVAILABLE,
            'pptx_processing': PPTX_AVAILABLE,
            'xlsx_processing': PANDAS_AVAILABLE,
            'pdf_ocr': PDF2IMAGE_AVAILABLE,
            'yolo': self._check_yolo_availability(),
            'ocr': self._check_ocr_availability()
        }
        
        print(f"[📎 첨부파일 서비스 초기화] 사용 가능한 기능: {sum(self.features.values())}/{len(self.features)}")
        print(f"[📎 상세 기능 체크]")
        print(f"  - PIL: {PIL_AVAILABLE}")
        print(f"  - YOLO: {self.features['yolo']}")
        print(f"  - OCR: {self.features['ocr']}")
    
    def _check_yolo_availability(self):
        """YOLO 모델 사용 가능 여부 확인 - 강화된 버전"""
        try:
            print("[🔍 YOLO 가용성 체크 시작]")
            
            if not hasattr(self.ai_models, 'yolo_model'):
                print("[❗YOLO] ai_models에 yolo_model 속성이 없음")
                return False
            
            if self.ai_models.yolo_model is None:
                print("[❗YOLO] yolo_model이 None - 자동 로딩 시도")
                if self.ai_models.load_yolo_model():
                    print("[✅ YOLO] 자동 로딩 성공")
                    return True
                else:
                    print("[❗YOLO] 자동 로딩 실패")
                    return False
            
            # 모델 테스트
            print("[🧪 YOLO] 모델 테스트 시작")
            test_image = np.zeros((320, 320, 3), dtype=np.uint8)
            results = self.ai_models.yolo_model(test_image, conf=0.1, verbose=False)
            print(f"[✅ YOLO] 테스트 성공 - 결과: {len(results)}개")
            return True
            
        except Exception as e:
            print(f"[❗YOLO] 체크 오류: {str(e)}")
            return False
    
    def _check_ocr_availability(self):
        """OCR 모델 사용 가능 여부 확인"""
        try:
            if not hasattr(self.ai_models, 'ocr_reader'):
                return False
            
            if self.ai_models.ocr_reader is None:
                if self.ai_models.load_ocr_model():
                    return True
                else:
                    return False
            
            return True
            
        except Exception as e:
            print(f"[❗OCR] 체크 오류: {str(e)}")
            return False
    
    def process_email_attachments(self, email_message, email_subject, email_id):
        """이메일에서 첨부파일을 추출하고 처리 (캐싱 포함)"""
        cache_key = f"email_{email_id}"
        
        # 캐시 확인
        if cache_key in self.attachment_cache:
            print(f"[📎 캐시 사용] {email_subject[:30]}...")
            return self.attachment_cache[cache_key]
        
        attachments = []
        print(f"[📎 새로운 첨부파일 처리] {email_subject[:30]}...")
        
        try:
            for part in email_message.walk():
                if part.get_content_disposition() == 'attachment':
                    attachment_info = self._process_single_attachment(part, email_subject)
                    if attachment_info:
                        attachments.append(attachment_info)
                        print(f"[📎 첨부파일 추가] {attachment_info.get('filename', 'Unknown')} - 타입: {attachment_info.get('type', 'Unknown')}")
        except Exception as e:
            print(f"[❗첨부파일 워킹 오류] {str(e)}")
        
        # 캐시 저장
        self.attachment_cache[cache_key] = attachments
        self._manage_cache_size()
        
        print(f"[✅ 첨부파일 처리 완료] {len(attachments)}개 처리됨")
        return attachments
    
    def _process_single_attachment(self, part, email_subject):
        """개별 첨부파일 처리 - 강화된 버전"""
        try:
            filename = self._decode_filename(part.get_filename())
            if not filename:
                print("[⚠️ 첨부파일] 파일명이 없음")
                return None
            
            attachment_data = part.get_payload(decode=True)
            if not attachment_data:
                print(f"[⚠️ 첨부파일] 데이터가 없음: {filename}")
                return None
            
            file_ext = Path(filename).suffix.lower()
            mime_type = part.get_content_type()
            
            print(f"[📄 첨부파일 분석] {filename} - 확장자: {file_ext}, MIME: {mime_type}, 크기: {len(attachment_data)} bytes")
            
            attachment_info = {
                'filename': filename,
                'size': len(attachment_data),
                'mime_type': mime_type,
                'extension': file_ext
            }
            
            # 파일 타입별 처리
            if file_ext in {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp'} or 'image' in mime_type:
                print(f"[🖼️ 이미지 파일 감지] {filename}")
                attachment_info.update(self._process_image(attachment_data, filename))
            elif file_ext == '.pdf' or 'pdf' in mime_type:
                print(f"[📄 PDF 파일 감지] {filename}")
                attachment_info.update(self._process_pdf(attachment_data, filename))
            elif file_ext == '.docx' or 'wordprocessingml' in mime_type:
                print(f"[📝 Word 파일 감지] {filename}")
                attachment_info.update(self._process_docx(attachment_data, filename))
            elif file_ext == '.pptx' or 'presentationml' in mime_type:
                print(f"[📊 PowerPoint 파일 감지] {filename}")
                attachment_info.update(self._process_pptx(attachment_data, filename))
            elif file_ext in ['.xlsx', '.xls'] or 'spreadsheetml' in mime_type:
                print(f"[📈 Excel 파일 감지] {filename}")
                attachment_info.update(self._process_xlsx(attachment_data, filename))
            else:
                print(f"[❓ 기타 파일] {filename}")
                attachment_info.update({'type': 'other', 'processing_method': 'metadata_only'})
            
            return attachment_info
            
        except Exception as e:
            print(f"[❗첨부파일 처리 오류] {filename if 'filename' in locals() else 'Unknown'}: {str(e)}")
            return None
    
    def _process_image(self, attachment_data, filename):
        """이미지 처리 (YOLO + OCR) - 완전 개선된 버전"""
        print(f"[🖼️ 이미지 처리 시작] {filename}")
        
        try:
            if not PIL_AVAILABLE:
                print(f"[❗이미지] PIL 사용 불가 - {filename}")
                return {'type': 'image', 'error': 'PIL not available', 'processing_method': 'disabled'}
            
            # YOLO 객체 인식
            yolo_detections = []
            yolo_success = False
            
            # YOLO 상태 재확인
            yolo_available = (self.features.get('yolo', False) and 
                            hasattr(self.ai_models, 'yolo_model') and 
                            self.ai_models.yolo_model is not None)
            
            print(f"[🔍 YOLO 상태 확인] 사용가능: {yolo_available} - {filename}")
            
            if yolo_available:
                print(f"[🤖 YOLO 처리 시작] {filename}")
                yolo_detections = self._yolo_detect_objects(attachment_data, filename)
                yolo_success = len(yolo_detections) > 0
                print(f"[🎯 YOLO 최종 결과] {len(yolo_detections)}개 객체 탐지 - {filename}")
            else:
                print(f"[⚠️ YOLO 스킵] 모델 사용 불가 - {filename}")
                print(f"    - features['yolo']: {self.features.get('yolo', False)}")
                print(f"    - yolo_model 존재: {hasattr(self.ai_models, 'yolo_model')}")
                print(f"    - yolo_model None 여부: {getattr(self.ai_models, 'yolo_model', None) is None}")
            
            # OCR 텍스트 추출
            ocr_result = {'text': '', 'success': False}
            ocr_available = (self.features.get('ocr', False) and 
                           hasattr(self.ai_models, 'ocr_reader') and 
                           self.ai_models.ocr_reader is not None)
            
            if ocr_available:
                print(f"[📝 OCR 처리 시작] {filename}")
                ocr_result = self._extract_text_with_ocr(attachment_data, filename)
                print(f"[📝 OCR 결과] 성공: {ocr_result.get('success')}, 텍스트 길이: {len(ocr_result.get('text', ''))} - {filename}")
            else:
                print(f"[⚠️ OCR 스킵] 모델 사용 불가 - {filename}")
            
            result = {
                'type': 'image',
                'yolo_detections': yolo_detections,
                'detected_objects': [det['class'] for det in yolo_detections],
                'object_count': len(yolo_detections),
                'extracted_text': ocr_result.get('text', ''),
                'ocr_success': ocr_result.get('success', False),
                'processing_method': f"YOLO({len(yolo_detections)}) + OCR({ocr_result.get('success', False)})",
                'yolo_success': yolo_success,
                'filename': filename
            }
            
            # 텍스트 요약 생성
            if ocr_result.get('success') and ocr_result.get('text'):
                result['text_summary'] = self._summarize_document(
                    ocr_result['text'], filename, 'image_with_text'
                )
            
            print(f"[✅ 이미지 처리 완료] {filename} - YOLO: {yolo_success} ({len(yolo_detections)}개), OCR: {ocr_result.get('success', False)}")
            return result
            
        except Exception as e:
            print(f"[❗이미지 처리 오류] {filename}: {str(e)}")
            import traceback
            traceback.print_exc()
            return {'type': 'image', 'error': str(e), 'processing_method': 'failed', 'filename': filename}
    
    def _yolo_detect_objects(self, image_data, filename):
        """YOLO 객체 인식 - 디버깅 강화"""
        try:
            print(f"[🔍 YOLO 추론 시작] {filename}")
            
            # 기본 체크
            if not PIL_AVAILABLE:
                print(f"[❗YOLO] PIL 사용 불가 - {filename}")
                return []
            
            if self.ai_models.yolo_model is None:
                print(f"[❗YOLO] 모델이 None - {filename}")
                return []
            
            # 이미지 로드
            print(f"[📁 이미지 로드] {filename}")
            image = Image.open(io.BytesIO(image_data))
            original_size = image.size
            print(f"[📏 이미지 크기] {original_size} - {filename}")
            
            # RGBA → RGB 변환
            if image.mode in ['RGBA', 'LA']:
                print(f"[🎨 색상 변환] {image.mode} → RGB - {filename}")
                rgb_image = Image.new('RGB', image.size, (255, 255, 255))
                rgb_image.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
                image = rgb_image
            elif image.mode != 'RGB':
                print(f"[🎨 색상 변환] {image.mode} → RGB - {filename}")
                image = image.convert('RGB')
            
            # NumPy 배열 변환
            image_np = np.array(image)
            print(f"[🔢 NumPy 변환] shape: {image_np.shape}, dtype: {image_np.dtype} - {filename}")
            
            # YOLO 모델 실행
            print(f"[🤖 YOLO 모델 실행 시작] - {filename}")
            try:
                results = self.ai_models.yolo_model(image_np, conf=0.1, verbose=False)
                print(f"[✅ YOLO 모델 실행 완료] 결과 개수: {len(results)} - {filename}")
            except Exception as model_error:
                print(f"[❗YOLO 모델 실행 실패] {str(model_error)} - {filename}")
                return []
            
            # 결과 처리
            detections = []
            if len(results) > 0:
                print(f"[📋 YOLO 결과 분석] 첫 번째 결과 타입: {type(results[0])} - {filename}")
                
                if hasattr(results[0], 'boxes') and results[0].boxes is not None:
                    boxes = results[0].boxes
                    print(f"[📦 탐지 박스] {len(boxes)}개 발견 - {filename}")
                    
                    for i in range(len(boxes)):
                        try:
                            conf = float(boxes.conf[i].cpu().numpy())
                            cls = int(boxes.cls[i].cpu().numpy())
                            class_name = self.ai_models.yolo_model.names[cls]
                            
                            detections.append({
                                'class': class_name,
                                'confidence': conf,
                                'class_id': cls
                            })
                            
                            print(f"[🎯 객체 탐지] {class_name} (신뢰도: {conf:.3f}) - {filename}")
                            
                        except Exception as box_error:
                            print(f"[❗박스 처리 오류] 인덱스 {i}: {str(box_error)} - {filename}")
                            continue
                else:
                    print(f"[❌ 박스 없음] boxes 속성이 None이거나 없음 - {filename}")
            else:
                print(f"[❌ 결과 없음] YOLO 결과가 비어있음 - {filename}")
            
            # 신뢰도순 정렬
            detections.sort(key=lambda x: x['confidence'], reverse=True)
            
            print(f"[✅ YOLO 완료] {len(detections)}개 객체 최종 탐지 - {filename}")
            for det in detections:
                print(f"  - {det['class']}: {det['confidence']:.3f}")
            
            return detections
            
        except Exception as e:
            print(f"[❗YOLO 처리 오류] {filename}: {str(e)}")
            import traceback
            traceback.print_exc()
            return []
    
    def _extract_text_with_ocr(self, attachment_data, filename):
        """OCR 텍스트 추출"""
        try:
            if not PIL_AVAILABLE:
                return {'text': '', 'success': False, 'error': 'PIL not available'}
            
            image = Image.open(io.BytesIO(attachment_data))
            
            # 이미지 전처리
            if image.mode in ['RGBA', 'LA']:
                rgb_image = Image.new('RGB', image.size, (255, 255, 255))
                rgb_image.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
                image = rgb_image
            elif image.mode != 'RGB':
                image = image.convert('RGB')
            
            image_np = np.array(image)
            
            # OCR 수행
            result = self.ai_models.ocr_reader.readtext(image_np, paragraph=True)
            
            text = ""
            for detection in result:
                if len(detection) >= 3:
                    text_content = detection[1]
                    confidence = detection[2]
                    if confidence > 0.5:
                        text += text_content + " "
            
            return {
                'text': text.strip(),
                'success': bool(text.strip()),
                'method': 'ocr'
            }
            
        except Exception as e:
            print(f"[❗OCR 오류] {str(e)}")
            return {'text': '', 'success': False, 'error': str(e)}
    
    def _decode_filename(self, filename):
        """파일명 디코딩"""
        if not filename:
            return None
        
        try:
            from email.header import decode_header
            decoded_parts = decode_header(filename)
            if decoded_parts and decoded_parts[0]:
                decoded_filename = decoded_parts[0]
                if isinstance(decoded_filename[0], bytes):
                    return decoded_filename[0].decode(decoded_filename[1] or 'utf-8')
                else:
                    return decoded_filename[0]
        except:
            pass
        
        return filename
    
    def _manage_cache_size(self):
        """캐시 크기 관리"""
        if len(self.attachment_cache) > self.config.MAX_CACHE_SIZE:
            oldest_key = next(iter(self.attachment_cache))
            del self.attachment_cache[oldest_key]
            print(f"[🗑️ 캐시 정리] 오래된 항목 삭제: {oldest_key}")
    
    def generate_attachment_summary(self, attachments):
        """첨부파일 요약 생성"""
        if not attachments:
            return ""
        
        total_files = len(attachments)
        
        # 파일 타입별 분류
        images = [att for att in attachments if att.get('type') == 'image']
        documents = [att for att in attachments if att.get('type', '').startswith('document_')]
        others = [att for att in attachments if att.get('type') not in ['image'] and not att.get('type', '').startswith('document_')]
        
        summary_parts = []
        
        if images:
            total_objects = sum(att.get('object_count', 0) for att in images)
            ocr_texts = [att for att in images if att.get('ocr_success')]
            
            if total_objects > 0:
                summary_parts.append(f"이미지 {len(images)}개({total_objects}개 객체)")
            else:
                summary_parts.append(f"이미지 {len(images)}개")
                
            if ocr_texts:
                summary_parts.append(f"텍스트 추출 {len(ocr_texts)}개")
        
        if documents:
            doc_types = {}
            successful_extractions = 0
            
            for doc in documents:
                doc_type = doc.get('type', '').replace('document_', '')
                doc_types[doc_type] = doc_types.get(doc_type, 0) + 1
                
                if doc.get('extraction_success'):
                    successful_extractions += 1
            
            for doc_type, count in doc_types.items():
                type_names = {
                    'pdf': 'PDF', 
                    'word': 'Word', 
                    'presentation': 'PPT', 
                    'spreadsheet': 'Excel'
                }
                type_name = type_names.get(doc_type, doc_type.upper())
                summary_parts.append(f"{type_name} {count}개")
            
            if successful_extractions > 0:
                summary_parts.append(f"요약 가능 {successful_extractions}개")
        
        if others:
            summary_parts.append(f"기타 {len(others)}개")
        
        if summary_parts:
            return f"📎 {total_files}개 파일: " + ", ".join(summary_parts)
        else:
            return f"📎 {total_files}개 파일"
    
    def clear_cache(self):
        """캐시 초기화"""
        cache_count = len(self.attachment_cache)
        self.attachment_cache.clear()
        return cache_count
    
    def get_available_features(self):
        """사용 가능한 기능 목록 반환"""
        return self.features
    
    def _process_pdf(self, attachment_data, filename):
        """PDF 파일 처리 - 텍스트 추출 및 OCR"""
        print(f"[📄 PDF 처리 시작] {filename}")
        
        try:
            extracted_text = ""
            extraction_method = ""
            extraction_success = False
            
            # 1. pdfplumber로 텍스트 추출 시도
            if PDFPLUMBER_AVAILABLE:
                print(f"[📖 pdfplumber 추출 시도] {filename}")
                try:
                    with pdfplumber.open(io.BytesIO(attachment_data)) as pdf:
                        pdf_text = ""
                        for page_num, page in enumerate(pdf.pages):
                            page_text = page.extract_text()
                            if page_text:
                                pdf_text += f"\n--- 페이지 {page_num + 1} ---\n{page_text}"
                        
                        if pdf_text.strip():
                            extracted_text = pdf_text.strip()
                            extraction_method = "pdfplumber"
                            extraction_success = True
                            print(f"[✅ pdfplumber 성공] {len(extracted_text)}자 추출 - {filename}")
                        else:
                            print(f"[⚠️ pdfplumber] 텍스트 없음 - {filename}")
                            
                except Exception as e:
                    print(f"[❗pdfplumber 실패] {str(e)} - {filename}")
            
            # 2. PyPDF2로 백업 시도 (pdfplumber 실패 시)
            if not extraction_success and PYPDF2_AVAILABLE:
                print(f"[📖 PyPDF2 백업 추출 시도] {filename}")
                try:
                    reader = PyPDF2.PdfReader(io.BytesIO(attachment_data))
                    pdf_text = ""
                    for page_num, page in enumerate(reader.pages):
                        page_text = page.extract_text()
                        if page_text:
                            pdf_text += f"\n--- 페이지 {page_num + 1} ---\n{page_text}"
                    
                    if pdf_text.strip():
                        extracted_text = pdf_text.strip()
                        extraction_method = "PyPDF2"
                        extraction_success = True
                        print(f"[✅ PyPDF2 성공] {len(extracted_text)}자 추출 - {filename}")
                    else:
                        print(f"[⚠️ PyPDF2] 텍스트 없음 - {filename}")
                        
                except Exception as e:
                    print(f"[❗PyPDF2 실패] {str(e)} - {filename}")
            
            # 3. OCR 시도 (텍스트 추출 실패하거나 결과가 부족한 경우)
            ocr_text = ""
            ocr_success = False
            
            if (not extraction_success or len(extracted_text) < 100) and PDF2IMAGE_AVAILABLE:
                print(f"[📷 PDF OCR 시도] {filename}")
                try:
                    # PDF를 이미지로 변환
                    images = convert_from_bytes(attachment_data, dpi=200)
                    print(f"[🖼️ PDF 변환] {len(images)}페이지 → 이미지 - {filename}")
                    
                    ocr_texts = []
                    for i, image in enumerate(images[:5]):  # 최대 5페이지만 OCR
                        if self.features.get('ocr') and hasattr(self.ai_models, 'ocr_reader') and self.ai_models.ocr_reader:
                            try:
                                image_np = np.array(image)
                                result = self.ai_models.ocr_reader.readtext(image_np, paragraph=True)
                                
                                page_text = ""
                                for detection in result:
                                    if len(detection) >= 3 and detection[2] > 0.5:
                                        page_text += detection[1] + " "
                                
                                if page_text.strip():
                                    ocr_texts.append(f"\n--- OCR 페이지 {i + 1} ---\n{page_text.strip()}")
                                    print(f"[📝 OCR 페이지 {i + 1}] {len(page_text)}자 추출 - {filename}")
                            
                            except Exception as ocr_error:
                                print(f"[❗OCR 페이지 {i + 1} 실패] {str(ocr_error)} - {filename}")
                                continue
                    
                    if ocr_texts:
                        ocr_text = "\n".join(ocr_texts)
                        ocr_success = True
                        print(f"[✅ PDF OCR 완료] {len(ocr_text)}자 총 추출 - {filename}")
                        
                        # 기존 텍스트가 부족하면 OCR 결과로 대체
                        if not extraction_success or len(extracted_text) < len(ocr_text):
                            extracted_text = ocr_text
                            extraction_method = "OCR"
                            extraction_success = True
                            
                except Exception as e:
                    print(f"[❗PDF OCR 실패] {str(e)} - {filename}")
            
            # 결과 구성
            result = {
                'type': 'document_pdf',
                'extracted_text': extracted_text,
                'text_length': len(extracted_text),
                'extraction_success': extraction_success,
                'extraction_method': extraction_method,
                'ocr_text': ocr_text,
                'ocr_success': ocr_success,
                'processing_method': f"PDF텍스트({extraction_method}) + OCR({ocr_success})",
                'filename': filename
            }
            
            # 텍스트 요약 생성
            if extraction_success and extracted_text:
                summary = self._summarize_document(extracted_text, filename, 'pdf_document')
                result['summary'] = summary
                result['text_summary'] = summary  # 호환성
                result['document_summary'] = summary  # API 응답용
            
            status = "성공" if extraction_success else "실패"
            method_info = f"{extraction_method} + OCR" if ocr_success else extraction_method
            print(f"[✅ PDF 처리 완료] {filename} - {status} ({method_info}), {len(extracted_text)}자")
            
            return result
            
        except Exception as e:
            print(f"[❗PDF 처리 오류] {filename}: {str(e)}")
            import traceback
            traceback.print_exc()
            return {
                'type': 'document_pdf',
                'error': str(e),
                'processing_method': 'failed',
                'extraction_success': False,
                'filename': filename
            }

    def _process_docx(self, attachment_data, filename):
        """Word 문서 처리"""
        print(f"[📝 Word 처리 시작] {filename}")
        
        try:
            if not DOCX_AVAILABLE:
                return {
                    'type': 'document_word',
                    'error': 'python-docx not available',
                    'processing_method': 'disabled',
                    'filename': filename
                }
            
            doc = Document(io.BytesIO(attachment_data))
            
            # 텍스트 추출
            full_text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    full_text.append(paragraph.text)
            
            # 표 내용 추출
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        full_text.append(" | ".join(row_text))
            
            extracted_text = "\n".join(full_text)
            extraction_success = bool(extracted_text.strip())
            
            result = {
                'type': 'document_word',
                'extracted_text': extracted_text,
                'text_length': len(extracted_text),
                'extraction_success': extraction_success,
                'extraction_method': 'python-docx',
                'processing_method': 'DOCX parser',
                'filename': filename
            }
            
            # 요약 생성
            if extraction_success:
                result['summary'] = self._summarize_document(
                    extracted_text, filename, 'word_document'
                )
                result['text_summary'] = result['summary']
            
            print(f"[✅ Word 처리 완료] {filename} - {len(extracted_text)}자 추출")
            return result
            
        except Exception as e:
            print(f"[❗Word 처리 오류] {filename}: {str(e)}")
            return {
                'type': 'document_word',
                'error': str(e),
                'processing_method': 'failed',
                'extraction_success': False,
                'filename': filename
            }

    def _process_pptx(self, attachment_data, filename):
        """PowerPoint 문서 처리"""
        print(f"[📊 PowerPoint 처리 시작] {filename}")
        
        try:
            if not PPTX_AVAILABLE:
                return {
                    'type': 'document_presentation',
                    'error': 'python-pptx not available',
                    'processing_method': 'disabled',
                    'filename': filename
                }
            
            prs = Presentation(io.BytesIO(attachment_data))
            
            # 텍스트 추출
            full_text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"\n--- 슬라이드 {slide_num} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # 헤더 외에 텍스트가 있는 경우
                    full_text.extend(slide_text)
            
            extracted_text = "\n".join(full_text)
            extraction_success = bool(extracted_text.strip())
            
            result = {
                'type': 'document_presentation',
                'extracted_text': extracted_text,
                'text_length': len(extracted_text),
                'slide_count': len(prs.slides),
                'extraction_success': extraction_success,
                'extraction_method': 'python-pptx',
                'processing_method': 'PPTX parser',
                'filename': filename
            }
            
            # 요약 생성
            if extraction_success:
                result['summary'] = self._summarize_document(
                    extracted_text, filename, 'presentation'
                )
                result['text_summary'] = result['summary']
            
            print(f"[✅ PowerPoint 처리 완료] {filename} - {len(prs.slides)}슬라이드, {len(extracted_text)}자 추출")
            return result
            
        except Exception as e:
            print(f"[❗PowerPoint 처리 오류] {filename}: {str(e)}")
            return {
                'type': 'document_presentation',
                'error': str(e),
                'processing_method': 'failed',
                'extraction_success': False,
                'filename': filename
            }

    def _process_xlsx(self, attachment_data, filename):
        """Excel 문서 처리"""
        print(f"[📈 Excel 처리 시작] {filename}")
        
        try:
            if not PANDAS_AVAILABLE:
                return {
                    'type': 'document_spreadsheet',
                    'error': 'pandas not available',
                    'processing_method': 'disabled',
                    'filename': filename
                }
            
            # Excel 파일 읽기
            excel_file = pd.ExcelFile(io.BytesIO(attachment_data))
            
            full_text = []
            sheet_info = []
            
            for sheet_name in excel_file.sheet_names:
                try:
                    df = pd.read_excel(excel_file, sheet_name=sheet_name)
                    
                    # 시트 정보 저장
                    sheet_info.append({
                        'name': sheet_name,
                        'rows': len(df),
                        'columns': len(df.columns)
                    })
                    
                    # 텍스트 변환
                    full_text.append(f"\n--- 시트: {sheet_name} ---")
                    
                    # 컬럼 헤더
                    if not df.empty:
                        full_text.append("컬럼: " + " | ".join(str(col) for col in df.columns))
                        
                        # 데이터 (최대 10행만)
                        for idx, row in df.head(10).iterrows():
                            row_data = []
                            for val in row:
                                if pd.notna(val):
                                    row_data.append(str(val))
                                else:
                                    row_data.append("")
                            full_text.append(" | ".join(row_data))
                        
                        if len(df) > 10:
                            full_text.append(f"... (총 {len(df)}행)")
                            
                except Exception as sheet_error:
                    print(f"[❗시트 처리 오류] {sheet_name}: {str(sheet_error)} - {filename}")
                    continue
            
            extracted_text = "\n".join(full_text)
            extraction_success = bool(extracted_text.strip())
            
            result = {
                'type': 'document_spreadsheet',
                'extracted_text': extracted_text,
                'text_length': len(extracted_text),
                'sheet_info': sheet_info,
                'sheet_count': len(excel_file.sheet_names),
                'extraction_success': extraction_success,
                'extraction_method': 'pandas',
                'processing_method': 'Excel parser',
                'filename': filename
            }
            
            # 요약 생성
            if extraction_success:
                result['summary'] = self._summarize_document(
                    extracted_text, filename, 'spreadsheet'
                )
                result['text_summary'] = result['summary']
            
            print(f"[✅ Excel 처리 완료] {filename} - {len(excel_file.sheet_names)}시트, {len(extracted_text)}자 추출")
            return result
            
        except Exception as e:
            print(f"[❗Excel 처리 오류] {filename}: {str(e)}")
            return {
                'type': 'document_spreadsheet',
                'error': str(e),
                'processing_method': 'failed',
                'extraction_success': False,
                'filename': filename
            }

    def _summarize_document(self, text, filename, document_type):
        """문서 요약 생성"""
        try:
            if not text or len(text.strip()) < 50:
                return "텍스트가 너무 짧아 요약할 수 없습니다."
            
            # 간단한 키워드 기반 요약 (실제로는 AI 모델 사용 권장)
            lines = text.split('\n')
            important_lines = []
            
            # 키워드가 포함된 라인 찾기
            keywords = ['문제', '답', '해답', '정답', '점수', '총점', '기말고사', '중간고사', 
                    '시험', '과제', '데이터베이스', 'database', 'sql', 'query']
            
            for line in lines:
                line = line.strip()
                if len(line) > 10 and any(keyword in line.lower() for keyword in keywords):
                    important_lines.append(line)
                    if len(important_lines) >= 5:  # 최대 5줄
                        break
            
            if important_lines:
                return f"{filename} 주요 내용:\n" + "\n".join(f"• {line[:100]}..." if len(line) > 100 else f"• {line}" for line in important_lines)
            else:
                # 키워드가 없으면 첫 3줄 요약
                summary_lines = [line.strip() for line in lines[:3] if line.strip()]
                if summary_lines:
                    return f"{filename} 시작 부분:\n" + "\n".join(f"• {line[:100]}..." if len(line) > 100 else f"• {line}" for line in summary_lines)
                else:
                    return f"{filename}: 문서 내용을 분석했지만 의미있는 요약을 생성할 수 없습니다."
                    
        except Exception as e:
            print(f"[❗요약 생성 오류] {filename}: {str(e)}")
            return f"{filename}: 요약 생성 중 오류가 발생했습니다."
