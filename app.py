from flask import Flask, request, jsonify, session
from email.header import decode_header
from email.utils import parseaddr, parsedate_to_datetime
from flask_cors import CORS
from datetime import datetime
import imaplib
import smtplib
import traceback
from email.mime.text import MIMEText
from transformers import pipeline
from nomic import embed
from sklearn.metrics.pairwise import cosine_similarity
from nomic import login
import os
import hashlib
import uuid
from huggingface_hub import InferenceClient
import torch
from transformers import AutoTokenizer, AutoModelForCausalLM
from pprint import pprint
import re
import email as email_module
import email

# ✅ YOLOv8 관련 import 추가
import cv2
import numpy as np
from pathlib import Path
import base64
from PIL import Image
import io

# 보고서 파일 처리용 라이브러리들
import pdfplumber  # PDF 텍스트 추출
import PyPDF2     # PDF 백업 처리
from docx import Document  # Word 문서 처리  
from pptx import Presentation  # PowerPoint 처리
import pandas as pd  # Excel 처리
import easyocr  # OCR 처리
from pdf2image import convert_from_bytes  # PDF → 이미지 변환
import mimetypes  # MIME 타입 감지
from pathlib import Path
import tempfile


# to 대쉬보드용 추가할 코드

import dateutil.parser
from datetime import datetime, timedelta
import re
import json
import time

from pathlib import Path  # 이 줄을 기존 import들과 함께 추가

# YOLOv8 설치 확인 및 로딩
from ultralytics import YOLO
print("[✅ YOLOv8 사용 가능]")

login(token="토큰")

# Hugging Face 토큰 설정
os.environ['HF_TOKEN'] = '토근'

candidate_labels = [
    "university.",
    "spam mail.",
    "company.",
    "security alert."
]

app = Flask(__name__)
CORS(app, supports_credentials=True)  # 세션 쿠키 지원
app.secret_key = 'your-secret-key-here'  # 세션 암호화용 키

# ✅ Qwen 모델 전역 변수 (한 번만 로딩하여 성능 최적화)
qwen_model = None
qwen_tokenizer = None

# ✅ YOLO 모델 전역 변수 추가
yolo_model = None

# OCR 모델 전역 변수
ocr_reader = None

# 사용자별 데이터 저장소 (실제 운영에서는 Redis나 데이터베이스 사용 권장)
user_sessions = {}

# ====================================
# app.py에 추가할 코드 (user_sessions = {} 바로 아래)
# ====================================

# 파일 기반 저장을 위한 디렉토리 생성
USER_DATA_DIR = Path("user_sessions")
USER_DATA_DIR.mkdir(exist_ok=True)
print(f"[📁 데이터 저장소] 생성: {USER_DATA_DIR}")

def get_user_file_path(user_email):
    """사용자별 데이터 파일 경로"""
    user_hash = hashlib.md5(user_email.encode()).hexdigest()[:16]
    return USER_DATA_DIR / f"user_{user_hash}.json"

def save_user_session_to_file(user_email):
    """현재 세션을 파일에 저장"""
    try:
        user_key = get_user_key(user_email)
        if user_key not in user_sessions:
            return False
            
        file_path = get_user_file_path(user_email)
        session_data = user_sessions[user_key]
        
        save_data = {
            'user_email': user_email,
            'extracted_todos': session_data.get('extracted_todos', []),
            'last_emails': session_data.get('last_emails', []),
            'last_update': datetime.now().isoformat(),
            'login_time': session_data.get('login_time'),
            'session_id': session_data.get('session_id')
        }
        
        with open(file_path, 'w', encoding='utf-8') as f:
            json.dump(save_data, f, ensure_ascii=False, indent=2)
        
        todos_count = len(session_data.get('extracted_todos', []))
        print(f"[💾 세션 저장] {user_email}: {todos_count}개 할일")
        return True
        
    except Exception as e:
        print(f"[❗저장 실패] {user_email}: {str(e)}")
        return False

def load_user_session_from_file(user_email):
    """파일에서 세션 데이터 로드"""
    try:
        file_path = get_user_file_path(user_email)
        
        if not file_path.exists():
            print(f"[📁 새 사용자] {user_email}")
            return None
        
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        if data.get('user_email') == user_email:
            todos_count = len(data.get('extracted_todos', []))
            print(f"[📂 세션 복원] {user_email}: {todos_count}개 할일")
            return data
            
        return None
        
    except Exception as e:
        print(f"[❗로드 실패] {user_email}: {str(e)}")
        return None

# ====================================
# 기존 함수들 수정
# ====================================

# 기존 clear_user_session 함수를 이것으로 교체
def clear_user_session(email):
    """특정 사용자의 세션 수정 - 파일은 유지"""
    user_key = get_user_key(email)
    if user_key in user_sessions:
        # 파일에 저장 후 메모리에서만 삭제
        save_user_session_to_file(email)
        del user_sessions[user_key]
        print(f"[🗑️ 세션 정리] {email} - 파일 저장 후 메모리 정리")

# 기존 login_user 함수를 이것으로 교체
@app.route('/api/login', methods=['POST'])
def login_user():
    """개선된 로그인 - 파일에서 데이터 복원"""
    try:
        data = request.get_json()
        email = data.get('email', '')
        
        if email:
            # 이전 세션 메모리에서 삭제 (파일은 유지)
            clear_user_session(email)
            
            # 파일에서 세션 복원 시도
            saved_data = load_user_session_from_file(email)
            
            # 새 세션 생성
            session_id = get_session_id()
            user_key = get_user_key(email)
            
            if saved_data:
                # 파일에서 복원
                user_sessions[user_key] = {
                    'email': email,
                    'session_id': session_id,
                    'extracted_todos': saved_data.get('extracted_todos', []),
                    'last_emails': saved_data.get('last_emails', []),
                    'login_time': datetime.now().isoformat()
                }
                
                todos_count = len(saved_data.get('extracted_todos', []))
                print(f"[🔑 로그인 + 복원] {email}: {todos_count}개 할일 복원")
                
                return jsonify({
                    'success': True,
                    'message': f'로그인 성공 - {todos_count}개 할일 복원됨',
                    'session_id': session_id,
                    'restored_todos': todos_count
                })
            else:
                # 새 세션 생성
                user_sessions[user_key] = {
                    'email': email,
                    'session_id': session_id,
                    'last_emails': [],
                    'extracted_todos': [],
                    'login_time': datetime.now().isoformat()
                }
                
                print(f"[🔑 새 로그인] {email}")
                
                return jsonify({
                    'success': True,
                    'message': '로그인 성공 - 새 세션',
                    'session_id': session_id,
                    'restored_todos': 0
                })
        else:
            return jsonify({'error': '이메일이 필요합니다.'}), 400
            
    except Exception as e:
        print(f"[❗로그인 실패] {str(e)}")
        return jsonify({'error': str(e)}), 500

# 기존 todos_api_improved 함수에 자동 저장 추가
@app.route('/api/todos', methods=['GET', 'POST', 'PUT', 'DELETE'])
def todos_api_improved():
    """할일 API - 자동 파일 저장 추가"""
    try:
        if request.method == 'GET':
            user_email = request.args.get('email')
        else:
            user_email = request.json.get('email') if request.json else None
            
        if not user_email:
            return jsonify({"error": "이메일이 필요합니다."}), 400
        
        user_key = get_user_key(user_email)
        if user_key not in user_sessions:
            return jsonify({"error": "로그인이 필요합니다."}), 401
        
        if request.method == 'GET':
            todos = user_sessions[user_key].get('extracted_todos', [])
            return jsonify({
                "success": True,
                "todos": todos,
                "total_count": len(todos)
            })
        
        elif request.method == 'POST':
            # 새 할일 추가
            data = request.json
            
            existing_todos = user_sessions[user_key].get('extracted_todos', [])
            max_id = max([todo.get('id', 0) for todo in existing_todos] + [0])
            new_id = max_id + 1
            
            new_todo = {
                'id': new_id,
                'type': data.get('type', 'task'),
                'title': data.get('title', ''),
                'description': data.get('description', ''),
                'date': data.get('date'),
                'time': data.get('time'),
                'priority': data.get('priority', 'medium'),
                'status': 'pending',
                'editable_date': True,
                'source_email': {
                    'from': 'manual',
                    'subject': 'Manual Entry',
                    'date': datetime.now().isoformat(),
                    'type': 'manual_entry'
                }
            }
            
            existing_todos.append(new_todo)
            user_sessions[user_key]['extracted_todos'] = existing_todos
            
            # 파일에 자동 저장
            save_user_session_to_file(user_email)
            
            print(f"[✅ 할일 추가 + 저장] ID: {new_id}")
            
            return jsonify({
                "success": True,
                "todo": new_todo,
                "message": "할일이 추가되고 저장되었습니다."
            })
        
        elif request.method == 'PUT':
            # 할일 업데이트
            data = request.json
            todo_id = data.get('id')
            
            todos = user_sessions[user_key].get('extracted_todos', [])
            updated = False
            
            for todo in todos:
                if todo.get('id') == todo_id:
                    if 'status' in data:
                        todo['status'] = data['status']
                    if 'date' in data and todo.get('editable_date', True):
                        todo['date'] = data['date']
                    if 'time' in data and todo.get('editable_date', True):
                        todo['time'] = data['time']
                    
                    updated = True
                    break
            
            if updated:
                user_sessions[user_key]['extracted_todos'] = todos
                
                # 파일에 자동 저장
                save_user_session_to_file(user_email)
                
                print(f"[✅ 할일 업데이트 + 저장] ID: {todo_id}")
                
                return jsonify({
                    "success": True,
                    "message": "할일이 업데이트되고 저장되었습니다."
                })
            else:
                return jsonify({"error": "해당 할일을 찾을 수 없습니다."}), 404
        
        elif request.method == 'DELETE':
            # 할일 삭제
            data = request.json
            todo_id = data.get('id')
            
            todos = user_sessions[user_key].get('extracted_todos', [])
            original_count = len(todos)
            
            todos = [todo for todo in todos if todo.get('id') != todo_id]
            
            if len(todos) < original_count:
                user_sessions[user_key]['extracted_todos'] = todos
                
                # 파일에 자동 저장
                save_user_session_to_file(user_email)
                
                print(f"[✅ 할일 삭제 + 저장] ID: {todo_id}")
                
                return jsonify({
                    "success": True,
                    "message": "할일이 삭제되고 저장되었습니다."
                })
            else:
                return jsonify({"error": "해당 할일을 찾을 수 없습니다."}), 404
        
    except Exception as e:
        print(f"[❗할일 API 오류] {str(e)}")
        return jsonify({"error": str(e)}), 500

# 기존 extract_todos_api 함수에 자동 저장 추가
@app.route('/api/extract-todos', methods=['POST'])
def extract_todos_api():
    """할일 추출 API - 자동 파일 저장 추가"""
    try:
        data = request.get_json()
        user_email = data.get("email", "")
        app_password = data.get("app_password", "")
        email_ids = data.get("email_ids", [])
        
        print(f"[📋 할일 추출] 사용자: {user_email}")
        
        user_key = get_user_key(user_email)
        if user_key not in user_sessions:
            return jsonify({"error": "로그인이 필요합니다."}), 401
        
        last_emails = user_sessions[user_key].get('last_emails', [])
        
        all_todos = []
        processed_count = 0
        
        emails_to_process = last_emails
        if email_ids:
            emails_to_process = [email for email in last_emails if email.get('id') in email_ids]
        
        for email_data in emails_to_process:
            try:
                result = extract_todos_from_email_improved(
                    email_body=email_data.get('body', ''),
                    email_subject=email_data.get('subject', ''),
                    email_from=email_data.get('from', ''),
                    email_date=email_data.get('date', '')
                )
                
                if result['success']:
                    for todo in result['todos']:
                        todo['source_email']['id'] = email_data.get('id')
                        todo['source_email']['subject'] = email_data.get('subject', '')
                    
                    all_todos.extend(result['todos'])
                    processed_count += 1
                    
            except Exception as e:
                continue
        
        # 기존 할일과 병합
        existing_todos = user_sessions[user_key].get('extracted_todos', [])
        existing_ids = {todo.get('id') for todo in existing_todos}
        
        new_todos = [todo for todo in all_todos if todo.get('id') not in existing_ids]
        final_todos = existing_todos + new_todos
        
        final_todos.sort(key=lambda x: x['date'] or '9999-12-31')
        
        user_sessions[user_key]['extracted_todos'] = final_todos
        
        # 파일에 자동 저장
        save_user_session_to_file(user_email)
        
        print(f"[✅ 할일 추출 + 저장] 총 {len(final_todos)}개 (신규 {len(new_todos)}개)")
        
        return jsonify({
            "success": True,
            "todos": final_todos,
            "total_count": len(final_todos),
            "new_todos": len(new_todos),
            "processed_emails": processed_count
        })
        
    except Exception as e:
        print(f"[❗할일 추출 오류] {str(e)}")
        return jsonify({"error": str(e)}), 500

# ✅ 이 줄을 추가하세요
attachment_cache = {}  # 첨부파일 처리 결과 캐시

# ✅ 첨부파일 저장 디렉토리 생성
ATTACHMENT_FOLDER = "static/attachments"
os.makedirs(ATTACHMENT_FOLDER, exist_ok=True)

# ===== YOLOv8 관련 함수들 추가 =====
def load_yolo_model():
    """YOLO 모델을 로딩하는 함수"""
    global yolo_model

    if yolo_model is None:
        try:
            print("[🤖 YOLOv8 모델 로딩 시작]")
            yolo_model = YOLO('yolov8n.pt')  # 경량 모델 사용
            print("[✅ YOLOv8 모델 로딩 완료]")
            return True
        except Exception as e:
            print(f"[❗YOLO 모델 로딩 실패] {str(e)}")
            return False
    return True
def process_image_with_yolo(image_data, confidence_threshold=0.2):
    """이미지 데이터를 YOLO로 처리 (PNG RGBA 문제 해결)"""
    global yolo_model
    
    if not load_yolo_model():
        return []
    
    try:
        # 바이너리 데이터를 이미지로 변환
        image = Image.open(io.BytesIO(image_data))
        
        # ✅ PNG RGBA → RGB 변환 (핵심 수정)
        if image.mode == 'RGBA' or image.mode == 'LA':
            print(f"[🔄 이미지 변환] {image.mode} → RGB")
            # 흰색 배경으로 RGBA → RGB 변환
            rgb_image = Image.new('RGB', image.size, (255, 255, 255))
            rgb_image.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
            image = rgb_image
        elif image.mode != 'RGB':
            print(f"[🔄 이미지 변환] {image.mode} → RGB")
            image = image.convert('RGB')
        
        image_np = np.array(image)
        print(f"[📐 이미지 크기] {image_np.shape}")  # 디버그용
        
        # YOLO 추론
        results = yolo_model(image_np, conf=confidence_threshold)
        
        detections = []
        if len(results) > 0 and results[0].boxes is not None:
            boxes = results[0].boxes
            for i in range(len(boxes)):
                conf = float(boxes.conf[i].cpu().numpy())
                cls = int(boxes.cls[i].cpu().numpy())
                class_name = yolo_model.names[cls]
                
                detections.append({
                    'class': class_name,
                    'confidence': conf,
                    'class_id': cls
                })
                
                print(f"  - {class_name}: {conf:.2f}")  # 디버그용
        
        return detections
        
    except Exception as e:
        print(f"[❗YOLO 이미지 처리 오류] {str(e)}")
        return []

def load_ocr_model():
    """EasyOCR 모델을 로딩하는 함수"""
    global ocr_reader
    
    if ocr_reader is None:
        try:
            print("[📖 EasyOCR 모델 로딩 시작]")
            ocr_reader = easyocr.Reader(['ko', 'en'])  # 한국어, 영어 지원
            print("[✅ EasyOCR 모델 로딩 완료]")
            return True
        except Exception as e:
            print(f"[❗EasyOCR 모델 로딩 실패] {str(e)}")
            return False
    return True

def extract_text_from_pdf(attachment_data):
    """PDF 파일에서 텍스트 추출 (OCR 백업 포함)"""
    try:
        # 1단계: pdfplumber로 텍스트 추출 시도
        print("[📄 PDF 텍스트 추출 시도]")
        
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_file:
            temp_file.write(attachment_data)
            temp_file_path = temp_file.name
        
        try:
            with pdfplumber.open(temp_file_path) as pdf:
                text = ""
                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n=== 페이지 {page_num + 1} ===\n{page_text}\n"
                
                if text.strip():
                    print(f"[✅ PDF 텍스트 추출 성공] {len(text)}자")
                    return {
                        'text': text.strip(),
                        'method': 'direct_extraction',
                        'pages': len(pdf.pages),
                        'success': True
                    }
        except Exception as e:
            print(f"[⚠️ pdfplumber 실패] {str(e)}")
        
        # 2단계: PyPDF2로 재시도
        try:
            with open(temp_file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n=== 페이지 {page_num + 1} ===\n{page_text}\n"
                
                if text.strip():
                    print(f"[✅ PyPDF2 텍스트 추출 성공] {len(text)}자")
                    return {
                        'text': text.strip(),
                        'method': 'pypdf2_extraction',
                        'pages': len(pdf_reader.pages),
                        'success': True
                    }
        except Exception as e:
            print(f"[⚠️ PyPDF2 실패] {str(e)}")
        
        # 3단계: OCR 백업 처리
        print("[🔍 PDF OCR 백업 처리 시작]")
        return extract_text_from_pdf_ocr(attachment_data)
        
    except Exception as e:
        print(f"[❗PDF 처리 실패] {str(e)}")
        return {
            'text': '',
            'method': 'failed',
            'error': str(e),
            'success': False
        }
    finally:
        # 임시 파일 정리
        try:
            os.unlink(temp_file_path)
        except:
            pass

def extract_text_from_pdf_ocr(attachment_data):
    """PDF를 이미지로 변환 후 OCR 처리"""
    try:
        if not load_ocr_model():
            return {'text': '', 'method': 'ocr_failed', 'success': False}
        
        # PDF를 이미지로 변환
        images = convert_from_bytes(attachment_data, dpi=200)
        
        all_text = ""
        for page_num, image in enumerate(images):
            # PIL Image를 numpy array로 변환
            image_np = np.array(image)
            
            # OCR 수행
            result = ocr_reader.readtext(image_np, paragraph=True)
            
            page_text = ""
            for detection in result:
                text = detection[1]
                confidence = detection[2]
                if confidence > 0.5:  # 신뢰도 50% 이상만
                    page_text += text + " "
            
            if page_text.strip():
                all_text += f"\n=== 페이지 {page_num + 1} (OCR) ===\n{page_text.strip()}\n"
        
        if all_text.strip():
            print(f"[✅ PDF OCR 성공] {len(all_text)}자")
            return {
                'text': all_text.strip(),
                'method': 'ocr_extraction',
                'pages': len(images),
                'success': True
            }
        else:
            return {
                'text': '',
                'method': 'ocr_no_text',
                'success': False
            }
            
    except Exception as e:
        print(f"[❗PDF OCR 실패] {str(e)}")
        return {
            'text': '',
            'method': 'ocr_failed',
            'error': str(e),
            'success': False
        }

def extract_text_from_docx(attachment_data):
    """Word 문서에서 텍스트 추출"""
    try:
        print("[📝 DOCX 텍스트 추출 시도]")
        
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as temp_file:
            temp_file.write(attachment_data)
            temp_file_path = temp_file.name
        
        try:
            doc = Document(temp_file_path)
            
            text = ""
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            # 표(Table) 내용도 추출
            for table in doc.tables:
                text += "\n=== 표 데이터 ===\n"
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text += " | ".join(row_text) + "\n"
            
            if text.strip():
                print(f"[✅ DOCX 텍스트 추출 성공] {len(text)}자")
                return {
                    'text': text.strip(),
                    'method': 'docx_extraction',
                    'paragraphs': len(doc.paragraphs),
                    'tables': len(doc.tables),
                    'success': True
                }
            else:
                return {
                    'text': '',
                    'method': 'docx_no_text',
                    'success': False
                }
                
        except Exception as e:
            print(f"[❗DOCX 처리 실패] {str(e)}")
            return {
                'text': '',
                'method': 'docx_failed',
                'error': str(e),
                'success': False
            }
        finally:
            try:
                os.unlink(temp_file_path)
            except:
                pass
                
    except Exception as e:
        print(f"[❗DOCX 파일 처리 실패] {str(e)}")
        return {
            'text': '',
            'method': 'docx_failed',
            'error': str(e),
            'success': False
        }

def extract_text_from_pptx(attachment_data):
    """PowerPoint 파일에서 텍스트 추출"""
    try:
        print("[📊 PPTX 텍스트 추출 시도]")
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            temp_file.write(attachment_data)
            temp_file_path = temp_file.name
        
        try:
            prs = Presentation(temp_file_path)
            
            text = ""
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n=== 슬라이드 {slide_num + 1} ===\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                    
                    # 표가 있는 경우
                    if shape.has_table:
                        text += "\n--- 표 ---\n"
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text += " | ".join(row_text) + "\n"
            
            if text.strip():
                print(f"[✅ PPTX 텍스트 추출 성공] {len(text)}자, {len(prs.slides)}개 슬라이드")
                return {
                    'text': text.strip(),
                    'method': 'pptx_extraction',
                    'slides': len(prs.slides),
                    'success': True
                }
            else:
                return {
                    'text': '',
                    'method': 'pptx_no_text',
                    'success': False
                }
                
        except Exception as e:
            print(f"[❗PPTX 처리 실패] {str(e)}")
            return {
                'text': '',
                'method': 'pptx_failed',
                'error': str(e),
                'success': False
            }
        finally:
            try:
                os.unlink(temp_file_path)
            except:
                pass
                
    except Exception as e:
        print(f"[❗PPTX 파일 처리 실패] {str(e)}")
        return {
            'text': '',
            'method': 'pptx_failed',
            'error': str(e),
            'success': False
        }

def extract_text_from_xlsx(attachment_data):
    """Excel 파일에서 데이터 추출"""
    try:
        print("[📊 XLSX 데이터 추출 시도]")
        
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as temp_file:
            temp_file.write(attachment_data)
            temp_file_path = temp_file.name
        
        try:
            # 모든 시트 읽기
            xl_file = pd.ExcelFile(temp_file_path)
            
            text = ""
            total_rows = 0
            
            for sheet_name in xl_file.sheet_names:
                df = pd.read_excel(temp_file_path, sheet_name=sheet_name)
                
                if not df.empty:
                    text += f"\n=== 시트: {sheet_name} ===\n"
                    
                    # 컬럼명 추가
                    text += "컬럼: " + " | ".join(str(col) for col in df.columns) + "\n\n"
                    
                    # 데이터 추가 (처음 20행만)
                    for idx, row in df.head(20).iterrows():
                        row_text = []
                        for value in row:
                            if pd.notna(value):
                                row_text.append(str(value))
                            else:
                                row_text.append("")
                        text += " | ".join(row_text) + "\n"
                    
                    total_rows += len(df)
                    
                    if len(df) > 20:
                        text += f"... (총 {len(df)}행 중 처음 20행만 표시)\n"
            
            if text.strip():
                print(f"[✅ XLSX 데이터 추출 성공] {len(text)}자, {total_rows}행")
                return {
                    'text': text.strip(),
                    'method': 'xlsx_extraction',
                    'sheets': len(xl_file.sheet_names),
                    'total_rows': total_rows,
                    'success': True
                }
            else:
                return {
                    'text': '',
                    'method': 'xlsx_no_data',
                    'success': False
                }
                
        except Exception as e:
            print(f"[❗XLSX 처리 실패] {str(e)}")
            return {
                'text': '',
                'method': 'xlsx_failed',
                'error': str(e),
                'success': False
            }
        finally:
            try:
                os.unlink(temp_file_path)
            except:
                pass
                
    except Exception as e:
        print(f"[❗XLSX 파일 처리 실패] {str(e)}")
        return {
            'text': '',
            'method': 'xlsx_failed',
            'error': str(e),
            'success': False
        }
def extract_text_with_ocr(attachment_data, filename):
    """일반 이미지 파일 OCR 처리 (오류 처리 강화)"""
    try:
        if not load_ocr_model():
            return {'text': '', 'method': 'ocr_model_failed', 'success': False}
        
        print(f"[🔍 이미지 OCR 처리] {filename}")
        
        # 이미지 로드 및 변환
        image = Image.open(io.BytesIO(attachment_data))
        
        # ✅ PNG RGBA → RGB 변환 (YOLO와 동일)
        if image.mode == 'RGBA' or image.mode == 'LA':
            print(f"[🔄 OCR 이미지 변환] {image.mode} → RGB")
            rgb_image = Image.new('RGB', image.size, (255, 255, 255))
            rgb_image.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
            image = rgb_image
        elif image.mode != 'RGB':
            print(f"[🔄 OCR 이미지 변환] {image.mode} → RGB")
            image = image.convert('RGB')
        
        image_np = np.array(image)
        
        # OCR 수행
        result = ocr_reader.readtext(image_np, paragraph=True)
        
        # ✅ OCR 결과 안전 처리
        text = ""
        if result and len(result) > 0:
            for detection in result:
                try:
                    # EasyOCR 결과 구조: [bbox, text, confidence]
                    if len(detection) >= 3:
                        text_content = detection[1]
                        confidence = detection[2]
                        if confidence > 0.5:  # 신뢰도 50% 이상만
                            text += text_content + " "
                except Exception as detail_error:
                    print(f"[⚠️ OCR 개별 결과 처리 오류] {str(detail_error)}")
                    continue
        
        if text.strip():
            print(f"[✅ 이미지 OCR 성공] {len(text)}자")
            return {
                'text': text.strip(),
                'method': 'image_ocr',
                'success': True
            }
        else:
            print(f"[📝 OCR 텍스트 없음] {filename}")
            return {
                'text': '',
                'method': 'ocr_no_text',
                'success': False
            }
            
    except Exception as e:
        print(f"[❗이미지 OCR 실패] {str(e)}")
        return {
            'text': '',
            'method': 'ocr_failed',
            'error': str(e),
            'success': False
        }
    
def summarize_document_with_llm(text, filename, file_type):
    """LLM을 이용한 문서 요약"""
    try:
        # 텍스트가 너무 길면 잘라내기 (4000자 제한)
        if len(text) > 4000:
            text = text[:4000] + "..."
        
        # Hugging Face 토큰 확인
        hf_token = os.getenv("HF_TOKEN")
        if not hf_token:
            print("[⚠️ HF_TOKEN 없음 - 간단 요약 사용]")
            return text[:300] + "..." if len(text) > 300 else text
        
        try:
            client = InferenceClient(
                model="Qwen/Qwen2.5-7B-Instruct",
                token=hf_token
            )
            
            prompt = f"""다음은 '{filename}' 파일의 내용입니다. 이 문서를 요약해주세요.

파일 형식: {file_type}
내용:
{text}

요약 지침:
1. 주요 내용을 3-5개 포인트로 요약
2. 핵심 키워드와 수치 포함
3. 150자 이내로 간결하게
4. 한국어로 응답

요약:"""
            
            messages = [
                {"role": "system", "content": "당신은 문서 요약 전문가입니다. 주어진 문서의 핵심 내용을 간결하고 정확하게 요약합니다."},
                {"role": "user", "content": prompt}
            ]
            
            response = client.chat_completion(
                messages=messages,
                max_tokens=200,
                temperature=0.3
            )
            
            summary = response.choices[0].message.content.strip()
            
            print(f"[✅ LLM 요약 완료] {filename} -> {len(summary)}자")
            return summary
            
        except Exception as e:
            print(f"[⚠️ LLM 요약 실패] {str(e)}")
            # 간단한 요약으로 fallback
            sentences = text.split('.')
            important_sentences = [s.strip() for s in sentences[:3] if len(s.strip()) > 10]
            return '. '.join(important_sentences) + '.' if important_sentences else text[:200] + "..."
            
    except Exception as e:
        print(f"[❗요약 처리 실패] {str(e)}")
        return text[:200] + "..." if len(text) > 200 else text

# ===== 4. 기존 extract_and_process_attachments 함수 확장 =====
# ===== 4. 기존 extract_and_process_attachments 함수 확장 =====

def extract_and_process_attachments_enhanced(email_message, email_subject, email_id):
    """이메일에서 첨부파일을 추출하고 YOLO + 보고서 처리 (캐싱 추가)"""
    global attachment_cache
    
    # ✅ 캐시 키 생성 (이메일 ID + 첨부파일 해시)
    cache_key = f"email_{email_id}"
    
    # ✅ 캐시에서 먼저 확인
    if cache_key in attachment_cache:
        print(f"[📎 캐시 사용] {email_subject[:30]}... - 첨부파일 처리 생략")
        return attachment_cache[cache_key]
    
    attachments = []
    print(f"[📎 새로운 첨부파일 처리] {email_subject[:30]}...")
    
    for part in email_message.walk():
        if part.get_content_disposition() == 'attachment':
            filename = part.get_filename()
            
            if filename:
                # MIME 디코딩
                try:
                    decoded_parts = decode_header(filename)
                    if decoded_parts and decoded_parts[0]:
                        decoded_filename = decoded_parts[0]
                        if isinstance(decoded_filename[0], bytes):
                            filename = decoded_filename[0].decode(decoded_filename[1] or 'utf-8')
                        else:
                            filename = decoded_filename[0]
                except:
                    pass
                
                attachment_data = part.get_payload(decode=True)
                if not attachment_data:
                    continue
                
                # ✅ 개별 첨부파일도 캐싱
                file_hash = hashlib.md5(attachment_data).hexdigest()[:8]
                file_cache_key = f"file_{filename}_{file_hash}"
                
                if file_cache_key in attachment_cache:
                    print(f"[📎 파일 캐시 사용] {filename}")
                    attachments.append(attachment_cache[file_cache_key])
                    continue
                
                # 파일 확장자 및 MIME 타입 확인
                file_ext = Path(filename).suffix.lower()
                mime_type = part.get_content_type()
                
                print(f"[📎 첨부파일 분석] {filename} ({file_ext}, {mime_type})")
                
                # 파일 형식별 처리
                attachment_info = {
                    'filename': filename,
                    'size': len(attachment_data),
                    'mime_type': mime_type,
                    'extension': file_ext
                }
                
                try:
                    # 1. 이미지 파일 처리 (기존 YOLO)
                    image_extensions = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp'}
                    if file_ext in image_extensions:
                        # YOLO 처리
                        yolo_detections = process_image_with_yolo(attachment_data)
                        
                        # 이미지에서 텍스트도 추출 (OCR) - 한 번만 실행
                        ocr_result = extract_text_with_ocr(attachment_data, filename)
                        
                        attachment_info.update({
                            'type': 'image',
                            'yolo_detections': yolo_detections,
                            'detected_objects': [det['class'] for det in yolo_detections],
                            'object_count': len(yolo_detections),
                            'extracted_text': ocr_result.get('text', ''),
                            'ocr_success': ocr_result.get('success', False),
                            'processing_method': f"YOLO + {'OCR' if ocr_result.get('success') else 'No OCR'}"
                        })
                        
                        if ocr_result.get('success') and ocr_result.get('text'):
                            summary = summarize_document_with_llm(
                                ocr_result['text'], filename, 'image_with_text'
                            )
                            attachment_info['text_summary'] = summary
                    
                    # 2. PDF 파일 처리
                    elif file_ext == '.pdf' or 'pdf' in mime_type:
                        pdf_result = extract_text_from_pdf(attachment_data)
                        
                        attachment_info.update({
                            'type': 'document_pdf',
                            'extracted_text': pdf_result.get('text', ''),
                            'extraction_success': pdf_result.get('success', False),
                            'extraction_method': pdf_result.get('method', 'unknown'),
                            'pages': pdf_result.get('pages', 0)
                        })
                        
                        if pdf_result.get('success') and pdf_result.get('text'):
                            summary = summarize_document_with_llm(
                                pdf_result['text'], filename, 'PDF 보고서'
                            )
                            attachment_info['document_summary'] = summary
                    
                    # 3. Word 문서 처리
                    elif file_ext == '.docx' or 'wordprocessingml' in mime_type:
                        docx_result = extract_text_from_docx(attachment_data)
                        
                        attachment_info.update({
                            'type': 'document_word',
                            'extracted_text': docx_result.get('text', ''),
                            'extraction_success': docx_result.get('success', False),
                            'paragraphs': docx_result.get('paragraphs', 0),
                            'tables': docx_result.get('tables', 0)
                        })
                        
                        if docx_result.get('success') and docx_result.get('text'):
                            summary = summarize_document_with_llm(
                                docx_result['text'], filename, 'Word 문서'
                            )
                            attachment_info['document_summary'] = summary
                    
                    # 4. PowerPoint 처리
                    elif file_ext == '.pptx' or 'presentationml' in mime_type:
                        pptx_result = extract_text_from_pptx(attachment_data)
                        
                        attachment_info.update({
                            'type': 'document_presentation',
                            'extracted_text': pptx_result.get('text', ''),
                            'extraction_success': pptx_result.get('success', False),
                            'slides': pptx_result.get('slides', 0)
                        })
                        
                        if pptx_result.get('success') and pptx_result.get('text'):
                            summary = summarize_document_with_llm(
                                pptx_result['text'], filename, 'PowerPoint 프레젠테이션'
                            )
                            attachment_info['document_summary'] = summary
                    
                    # 5. Excel 처리
                    elif file_ext in ['.xlsx', '.xls'] or 'spreadsheetml' in mime_type:
                        xlsx_result = extract_text_from_xlsx(attachment_data)
                        
                        attachment_info.update({
                            'type': 'document_spreadsheet',
                            'extracted_text': xlsx_result.get('text', ''),
                            'extraction_success': xlsx_result.get('success', False),
                            'sheets': xlsx_result.get('sheets', 0),
                            'total_rows': xlsx_result.get('total_rows', 0)
                        })
                        
                        if xlsx_result.get('success') and xlsx_result.get('text'):
                            summary = summarize_document_with_llm(
                                xlsx_result['text'], filename, 'Excel 스프레드시트'
                            )
                            attachment_info['document_summary'] = summary
                    
                    # 6. 기타 파일
                    else:
                        attachment_info.update({
                            'type': 'other',
                            'processing_method': 'metadata_only'
                        })
                    
                    # ✅ 개별 파일 캐시에 저장
                    attachment_cache[file_cache_key] = attachment_info
                    attachments.append(attachment_info)
                    
                    # 로그 출력
                    if attachment_info.get('extraction_success'):
                        print(f"[✅ 문서 처리 성공] {filename}: {attachment_info.get('type')}")
                    elif attachment_info.get('object_count', 0) > 0:
                        print(f"[✅ 이미지 처리 성공] {filename}: {attachment_info['object_count']}개 객체")
                    else:
                        print(f"[📎 파일 정보만 수집] {filename}")
                
                except Exception as e:
                    print(f"[❗첨부파일 처리 오류] {filename}: {str(e)}")
                    # 처리 실패해도 기본 정보는 저장
                    attachment_info.update({
                        'type': 'error',
                        'error': str(e),
                        'processing_method': 'failed'
                    })
                    attachments.append(attachment_info)
    
    # ✅ 전체 결과를 캐시에 저장
    attachment_cache[cache_key] = attachments
    
    # ✅ 캐시 크기 제한 (메모리 관리)
    if len(attachment_cache) > 100:  # 최대 100개 항목만 유지
        oldest_key = next(iter(attachment_cache))
        del attachment_cache[oldest_key]
        print(f"[🗑️ 캐시 정리] 오래된 항목 삭제: {oldest_key}")
    
    print(f"[✅ 첨부파일 처리 완료] {len(attachments)}개 처리됨 (캐시 저장)")
    return attachments
# app.py에 추가할 함수 (extract_and_process_attachments_enhanced 함수 다음에 추가)

def generate_enhanced_attachment_summary(attachments):
    """향상된 첨부파일 요약 생성"""
    if not attachments:
        return ""
    
    total_files = len(attachments)
    
    # 파일 타입별 분류
    images = [att for att in attachments if att.get('type') == 'image']
    documents = [att for att in attachments if att.get('type', '').startswith('document_')]
    others = [att for att in attachments if att.get('type') not in ['image'] and not att.get('type', '').startswith('document_')]
    
    summary_parts = []
    
    if images:
        total_objects = sum(att.get('object_count', 0) for att in images)
        ocr_texts = [att for att in images if att.get('ocr_success')]
        
        if total_objects > 0:
            summary_parts.append(f"이미지 {len(images)}개({total_objects}개 객체)")
        else:
            summary_parts.append(f"이미지 {len(images)}개")
            
        if ocr_texts:
            summary_parts.append(f"텍스트 추출 {len(ocr_texts)}개")
    
    if documents:
        # 문서 타입별 개수 계산
        doc_types = {}
        successful_extractions = 0
        
        for doc in documents:
            doc_type = doc.get('type', '').replace('document_', '')
            doc_types[doc_type] = doc_types.get(doc_type, 0) + 1
            
            if doc.get('extraction_success'):
                successful_extractions += 1
        
        # 문서 타입별 표시
        for doc_type, count in doc_types.items():
            type_names = {
                'pdf': 'PDF', 
                'word': 'Word', 
                'presentation': 'PPT', 
                'spreadsheet': 'Excel'
            }
            type_name = type_names.get(doc_type, doc_type.upper())
            summary_parts.append(f"{type_name} {count}개")
        
        # 성공적으로 처리된 문서 개수 표시
        if successful_extractions > 0:
            summary_parts.append(f"요약 가능 {successful_extractions}개")
    
    if others:
        summary_parts.append(f"기타 {len(others)}개")
    
    if summary_parts:
        return f"📎 {total_files}개 파일: " + ", ".join(summary_parts)
    else:
        return f"📎 {total_files}개 파일"

@app.route('/api/document-summary', methods=['POST'])
def document_summary():
    """특정 첨부파일의 상세 문서 요약 반환"""
    try:
        data = request.get_json()
        email_id = data.get("email_id")
        filename = data.get("filename", "")
        user_email = data.get("email", "")
        
        # 사용자 세션 확인
        user_key = get_user_key(user_email)
        if user_key not in user_sessions:
            return jsonify({"error": "로그인이 필요합니다."}), 401
        
        # 세션에서 해당 메일의 첨부파일 찾기
        last_emails = user_sessions[user_key].get('last_emails', [])
        target_attachment = None
        
        for email_data in last_emails:
            if email_data.get('id') == email_id:
                for attachment in email_data.get('attachments', []):
                    if attachment.get('filename') == filename:
                        target_attachment = attachment
                        break
                break
        
        if not target_attachment:
            return jsonify({"error": "해당 첨부파일을 찾을 수 없습니다."}), 404
        
        # 문서 요약 정보 반환
        response_data = {
            "success": True,
            "filename": filename,
            "file_type": target_attachment.get('type', 'unknown'),
            "size": target_attachment.get('size', 0),
            "extraction_success": target_attachment.get('extraction_success', False)
        }
        
        # 타입별 상세 정보 추가
        if target_attachment.get('type') == 'image':
            response_data.update({
                "yolo_detections": target_attachment.get('detected_objects', []),
                "object_count": target_attachment.get('object_count', 0),
                "ocr_text": target_attachment.get('extracted_text', ''),
                "text_summary": target_attachment.get('text_summary', '')
            })
        
        elif target_attachment.get('type', '').startswith('document_'):
            response_data.update({
                "extracted_text": target_attachment.get('extracted_text', '')[:1000],  # 처음 1000자만
                "document_summary": target_attachment.get('document_summary', ''),
                "extraction_method": target_attachment.get('extraction_method', ''),
                "full_text_available": len(target_attachment.get('extracted_text', '')) > 1000
            })
            
            # 파일 타입별 추가 정보
            if target_attachment.get('pages'):
                response_data['pages'] = target_attachment['pages']
            if target_attachment.get('slides'):
                response_data['slides'] = target_attachment['slides']
            if target_attachment.get('sheets'):
                response_data['sheets'] = target_attachment['sheets']
        
        return jsonify(response_data)
        
    except Exception as e:
        print(f"[❗문서 요약 API 오류] {str(e)}")
        return jsonify({"error": str(e)}), 500
    
# ===== 기존 Qwen 관련 함수들 =====
def load_qwen_model():
    """Qwen 모델을 로딩하는 함수"""
    global qwen_model, qwen_tokenizer
    
    if qwen_model is None:
        print("[🤖 Qwen 모델 로딩 시작]")
        try:
            device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
            model_id = "Qwen/Qwen1.5-1.8B-Chat"
            
            qwen_tokenizer = AutoTokenizer.from_pretrained(model_id, trust_remote_code=True)
            qwen_model = AutoModelForCausalLM.from_pretrained(
                model_id,
                device_map="auto",
                torch_dtype=torch.float16 if torch.cuda.is_available() else torch.float32,
                trust_remote_code=True
            )
            qwen_model.eval()
            print("[✅ Qwen 모델 로딩 완료]")
        except Exception as e:
            print(f"[❗Qwen 모델 로딩 실패] {str(e)}")
            # Qwen 로딩 실패해도 다른 기능은 정상 작동하도록

def extract_search_target_with_qwen(text):
    """Qwen을 이용하여 검색 대상 추출"""
    global qwen_model, qwen_tokenizer
    
    # 모델이 로딩되지 않았다면 로딩 시도
    if qwen_model is None:
        load_qwen_model()
    
    # 모델 로딩에 실패한 경우 간단한 키워드 추출로 fallback
    if qwen_model is None:
        print("[⚠️ Qwen 모델 없음 - 간단 추출 사용]")
        words = text.split()
        return " ".join(words[-2:]) if len(words) >= 2 else text
    
    try:
        prompt = (
            "<|im_start|>system\nYou are an email assistant. "
            "Your job is to extract the email address or name the user is referring to. "
            "You must always respond in the format: The user is referring to ... \n"
            "<|im_end|>\n"
            f"<|im_start|>user\n{text}<|im_end|>\n"
            "<|im_start|>assistant\n"
        )
        
        inputs = qwen_tokenizer(prompt, return_tensors="pt").to(qwen_model.device)
        
        with torch.no_grad():
            outputs = qwen_model.generate(
                **inputs,
                max_new_tokens=50,
                do_sample=False,
                eos_token_id=qwen_tokenizer.eos_token_id
            )
        
        decoded_output = qwen_tokenizer.decode(outputs[0], skip_special_tokens=True)
        print("\n Qwen 응답 전체:\n", decoded_output)
        # "assistant" 이후 텍스트만 가져옴
        if "assistant" in decoded_output:
            after_assistant = decoded_output.split("assistant")[-1].strip()
            prefix = "The user is referring to "
            if prefix in after_assistant:
                result = after_assistant.split(prefix)[-1].strip().rstrip(".").strip('"')
                return result
        

    except Exception as e:
        print(f"[⚠️ Qwen 추출 오류] {str(e)}")
        # 오류 시 간단한 키워드 추출로 fallback
        words = text.split()
        return " ".join(words[-2:]) if len(words) >= 2 else text

def search_emails_by_target(emails, search_target):
    """이메일 리스트에서 검색 대상으로 필터링"""
    results = []
    search_lower = search_target.lower()
    
    for mail in emails:
        # from 필드에서 검색
        if search_lower in mail["from"].lower():
            results.append(mail)
        # 제목에서도 검색
        elif search_lower in mail["subject"].lower():
            results.append(mail)
        # 이메일 주소만 추출해서 검색
        elif "@" in search_target:
            # 이메일 주소 패턴 매칭
            email_pattern = r'<([^>]+)>'
            email_match = re.search(email_pattern, mail["from"])
            if email_match and search_lower in email_match.group(1).lower():
                results.append(mail)
    
    return results

def get_session_id():
    """세션 ID 생성 또는 가져오기"""
    if 'session_id' not in session:
        session['session_id'] = str(uuid.uuid4())
    return session['session_id']

def get_user_key(email):
    """이메일 기반 사용자 키 생성"""
    return hashlib.md5(email.encode()).hexdigest()

# 요약 모델 로딩
summarizer = pipeline("summarization", model="facebook/bart-large-cnn", tokenizer="facebook/bart-large-cnn")

def build_ai_reply_prompt(sender, subject, body):
    """AI 답장을 위한 프롬프트 생성"""
    return f"""
You are a helpful email assistant that writes professional email replies.

Please read the following email and write a polite, professional reply in English:

---
From: {sender}
Subject: {subject}
Body: {body}
---

Instructions:
1. Identify the purpose of the email (invitation, question, information request, scheduling, etc.)
2. Write a concise (3-4 sentences), polite reply that directly addresses the purpose
3. Use a friendly yet professional tone
4. Only output the reply text (no analysis, no quotes, no original email content)

Reply:
""".strip()

# ===== API 엔드포인트들 =====

@app.route('/api/email-search', methods=['POST'])
def email_search():
    """이메일 검색 API - 변수명 충돌 해결"""
    try:
        data = request.get_json()
        user_input = data.get("user_input", "").strip()
        user_email = data.get("email", "")  # ✅ 변수명 변경
        app_password = data.get("app_password", "")
        
        print(f"[🔍 이메일 검색 요청] 사용자: {user_email}, 입력: {user_input}")
        
        if not all([user_input, user_email, app_password]):
            return jsonify({"error": "사용자 입력, 이메일, 앱 비밀번호가 모두 필요합니다."}), 400
        
        # 사용자 세션 확인
        user_key = get_user_key(user_email)
        if user_key not in user_sessions:
            return jsonify({"error": "로그인이 필요합니다."}), 401
        
        print("[🎯 실제 메일 검색 시작]")
        
        # Qwen을 이용해 검색 대상 추출
        try:
            search_target = extract_search_target_with_qwen(user_input)
            print(f"[🎯 검색 대상 추출] '{search_target}'")
        except Exception as e:
            print(f"[⚠️ Qwen 추출 실패] {str(e)}")
            # 간단한 이메일 추출 fallback
            import re
            email_pattern = r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
            emails_found = re.findall(email_pattern, user_input)
            if emails_found:
                search_target = emails_found[0]
            else:
                words = user_input.split()
                search_target = " ".join(words[-2:]) if len(words) >= 2 else user_input
        
        print(f"[🔍 최종 검색 대상] '{search_target}'")
        
        # 메일 서버 연결 및 검색
        try:
            mail = imaplib.IMAP4_SSL("imap.gmail.com")
            mail.login(user_email, app_password)  # ✅ 변수명 수정
            mail.select("inbox")
            print("[✅ 메일 서버 연결 성공]")
            
            # 검색 범위 설정
            N = 100
            status, data_result = mail.search(None, "ALL")
            all_mail_ids = data_result[0].split()
            mail_ids = all_mail_ids[-N:]  # 최근 N개
            
            print(f"[📊 검색 범위] 총 {len(all_mail_ids)}개 중 최근 {len(mail_ids)}개 검색")
            
            emails_found = []
            processed_count = 0
            
            for msg_id in mail_ids:
                try:
                    _, msg_data = mail.fetch(msg_id, "(RFC822)")
                    if not msg_data or not msg_data[0]:
                        continue
                        
                    # ✅ 올바른 모듈 사용
                    msg = email_module.message_from_bytes(msg_data[0][1])
                    processed_count += 1
                    
                    # 제목 디코딩
                    raw_subject = msg.get("Subject", "")
                    try:
                        decoded_parts = decode_header(raw_subject)
                        if decoded_parts and decoded_parts[0]:
                            decoded_subject = decoded_parts[0]
                            subject_bytes = decoded_subject[0]
                            subject_encoding = decoded_subject[1]
                            
                            if isinstance(subject_bytes, bytes):
                                if subject_encoding is None:
                                    subject_encoding = 'utf-8'
                                try:
                                    subject = subject_bytes.decode(subject_encoding)
                                except (UnicodeDecodeError, LookupError):
                                    for fallback_encoding in ['utf-8', 'latin-1', 'cp949', 'euc-kr']:
                                        try:
                                            subject = subject_bytes.decode(fallback_encoding)
                                            break
                                        except (UnicodeDecodeError, LookupError):
                                            continue
                                    else:
                                        subject = subject_bytes.decode('utf-8', errors='ignore')
                            else:
                                subject = str(subject_bytes)
                        else:
                            subject = "(제목 없음)"
                    except Exception as e:
                        subject = raw_subject if raw_subject else "(제목 없음)"
                    
                    # 발신자 정보
                    name, addr = parseaddr(msg.get("From"))
                    from_field = f"{name} <{addr}>" if name else addr
                    
                    # 날짜 처리
                    raw_date = msg.get("Date", "")
                    try:
                        date_obj = parsedate_to_datetime(raw_date)
                        date_str = date_obj.strftime("%Y-%m-%d %H:%M:%S")
                    except:
                        date_str = raw_date[:19] if len(raw_date) >= 19 else raw_date
                    
                    # 본문 추출
                    body = ""
                    try:
                        if msg.is_multipart():
                            for part in msg.walk():
                                if part.get_content_type() == "text/plain" and not part.get("Content-Disposition"):
                                    charset = part.get_content_charset() or "utf-8"
                                    body += part.get_payload(decode=True).decode(charset, errors="ignore")
                        else:
                            charset = msg.get_content_charset() or "utf-8"
                            body = msg.get_payload(decode=True).decode(charset, errors="ignore")
                        
                        body = body.strip()
                    except Exception as e:
                        body = ""
                    
                    # 검색 대상과 매칭 확인
                    search_in = f"{subject} {from_field} {body}".lower()
                    search_lower = search_target.lower()
                    
                    # 이메일 주소나 이름으로 검색
                    if (search_lower in search_in or 
                        any(part.strip() in search_in for part in search_lower.split() if part.strip())):
                        
                        emails_found.append({
                            "id": int(msg_id.decode()) if isinstance(msg_id, bytes) else int(msg_id),
                            "subject": subject,
                            "from": from_field,
                            "date": date_str,
                            "body": body[:500]  # 처음 500자만
                        })
                        
                        print(f"[✅ 매칭 발견] {from_field} -> {subject[:30]}...")
                        
                        if len(emails_found) >= 10:  # 최대 10개
                            break
                            
                except Exception as e:
                    print(f"[⚠️ 메일 처리 오류] {str(e)}")
                    continue
            
            mail.close()
            mail.logout()
            
            print(f"[📊 검색 완료] {processed_count}개 처리, {len(emails_found)}개 발견")
            
            return jsonify({
                "success": True,
                "search_target": search_target,
                "results": emails_found,
                "total_searched": processed_count,
                "found_count": len(emails_found),
                "confidence": 1.0,
                "detected_intent": "email_search_completed"
            })
            
        except Exception as e:
            print(f"[❗메일 서버 오류] {str(e)}")
            return jsonify({
                "success": False,
                "error": f"메일 서버 연결 오류: {str(e)}",
                "search_target": search_target if 'search_target' in locals() else user_input
            }), 500
            
    except Exception as e:
        print(f"[❗이메일 검색 오류] {str(e)}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route('/api/logout', methods=['POST'])
def logout_user():
    """사용자 로그아웃 - 세션 데이터 삭제"""
    try:
        data = request.get_json()
        email = data.get('email', '')
        
        if email:
            clear_user_session(email)
            session.clear()  # Flask 세션도 삭제
            
            return jsonify({
                'success': True,
                'message': '로그아웃 성공'
            })
        else:
            return jsonify({'error': '이메일이 필요합니다.'}), 400
            
    except Exception as e:
        print(f"[❗로그아웃 실패] {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/generate-ai-reply', methods=['POST'])
def generate_ai_reply():
    """AI 답장 생성 API"""
    try:
        data = request.get_json()
        sender = data.get('sender', '')
        subject = data.get('subject', '')
        body = data.get('body', '')
        current_user_email = data.get('email', '')  # 현재 사용자 이메일 추가
        
        print(f"[🤖 AI 답장 요청] User: {current_user_email}, From: {sender}, Subject: {subject[:50]}...")
        
        if not all([sender, subject, body, current_user_email]):
            return jsonify({'error': '발신자, 제목, 본문, 사용자 이메일이 모두 필요합니다.'}), 400
        
        # 사용자 세션 확인
        user_key = get_user_key(current_user_email)
        if user_key not in user_sessions:
            return jsonify({'error': '로그인이 필요합니다.'}), 401
        
        # Hugging Face 토큰 확인
        hf_token = os.getenv("HF_TOKEN")
        if not hf_token:
            return jsonify({'error': 'HF_TOKEN 환경 변수가 설정되어 있지 않습니다.'}), 500
        
        # InferenceClient 생성
        client = InferenceClient(
            model="Qwen/Qwen2.5-7B-Instruct",
            token=hf_token
        )
        
        # 프롬프트 생성
        user_prompt = build_ai_reply_prompt(sender, subject, body)
        
        # AI 답장 생성
        messages = [
            {"role": "system", "content": "You are a helpful email assistant that writes professional email replies."},
            {"role": "user", "content": user_prompt}
        ]
        
        response = client.chat_completion(
            messages=messages,
            max_tokens=256,
            temperature=0.7
        )
        
        # 답장 텍스트 추출
        ai_reply = response.choices[0].message.content.strip()
        
        print(f"[✅ AI 답장 생성 완료] User: {current_user_email}, 길이: {len(ai_reply)}자")
        
        return jsonify({
            'success': True,
            'ai_reply': ai_reply
        })
        
    except Exception as e:
        print(f"[❗AI 답장 생성 실패] {str(e)}")
        traceback.print_exc()
        return jsonify({'error': f'AI 답장 생성 실패: {str(e)}'}), 500

# ✅ 수정된 summary 함수 - 첨부파일 YOLO 처리 추가
@app.route('/api/summary', methods=['POST'])
def summary():
    try:
        data = request.get_json()
        username = data.get("email")
        app_password = data.get("app_password")

        # 사용자 키 생성 및 세션 확인
        user_key = get_user_key(username)
        
        print(f"[📧 메일 요청] 사용자: {username}")
        
        # 문자열 날짜를 datetime 객체로 변환
        after_date = data.get("after")
        after_dt = None
        if after_date:
            try:
                after_date_clean = after_date.replace("Z", "+00:00")
                after_dt = datetime.fromisoformat(after_date_clean)
                after_dt = after_dt.replace(tzinfo=None)
                print(f"[📅 필터링 기준] {after_dt} 이후 메일만 가져옴")
            except Exception as e:
                print("[⚠️ after_date 파싱 실패]", e)

        # 메일 서버 연결
        mail = imaplib.IMAP4_SSL("imap.gmail.com")
        mail.login(username, app_password)
        mail.select("inbox")

        # 메일 수 동적 결정
        if after_dt:
            N = 10
            print(f"[🔄 새로고침] 최근 {N}개 메일에서 {after_dt} 이후 메일 검색")
        else:
            N = 5
            print(f"[🆕 첫 로딩] 최근 {N}개 메일 가져옴")

        status, data = mail.search(None, "ALL")
        all_mail_ids = data[0].split()
        
        # 최신 메일부터 처리하도록 순서 수정
        mail_ids = all_mail_ids[-N:]
        mail_ids.reverse()

        emails = []
        processed_count = 0

        for msg_id in mail_ids:
            status, msg_data = mail.fetch(msg_id, "(RFC822)")
            if not msg_data or not msg_data[0]:
                continue

            raw_msg = msg_data[0][1]
            msg = email.message_from_bytes(raw_msg)

            # 제목 디코딩
            raw_subject = msg.get("Subject", "")
            decoded_parts = decode_header(raw_subject)
            if decoded_parts:
                decoded_subject = decoded_parts[0]
                subject = decoded_subject[0].decode(decoded_subject[1]) if isinstance(decoded_subject[0], bytes) else decoded_subject[0]
            else:
                subject = "(제목 없음)"

            # 보내는 사람
            name, addr = parseaddr(msg.get("From"))
            from_field = f"{name} <{addr}>" if name else addr

            # 날짜 처리
            raw_date = msg.get("Date", "")
            try:
                date_obj = parsedate_to_datetime(raw_date)
                date_obj = date_obj.replace(tzinfo=None)
                date_str = date_obj.strftime("%Y-%m-%d %H:%M:%S")
            except:
                date_obj = None
                date_str = raw_date[:19] if len(raw_date) >= 19 else raw_date

            # after_date 필터링
            if after_dt and date_obj:
                if date_obj <= after_dt:
                    print(f"[⏭️ 건너뛰기] {date_str} (기준: {after_dt})")
                    continue
                else:
                    print(f"[✅ 포함] {date_str} - {subject[:30]}...")

            # 본문 추출
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain" and not part.get("Content-Disposition"):
                        charset = part.get_content_charset() or "utf-8"
                        body += part.get_payload(decode=True).decode(charset, errors="ignore")
            else:
                charset = msg.get_content_charset() or "utf-8"
                body = msg.get_payload(decode=True).decode(charset, errors="ignore")

            body = body.strip()
            if not body:
                body = ""

            # ✅ 첨부파일 추출 및 YOLO 처리 추가
            email_id_str = msg_id.decode() if isinstance(msg_id, bytes) else str(msg_id)
            print(f"[📎 첨부파일 처리 시작] 메일: {subject[:30]}... (ID: {email_id_str})")
            
            # 첨부파일이 있는지 먼저 확인
            attachments = []
            try:
                print(f"[🔍 멀티파트 확인] {msg.is_multipart()}")
                
                if msg.is_multipart():
                    attachment_count = 0
                    for part in msg.walk():
                        content_disp = part.get_content_disposition()
                        if content_disp == 'attachment':
                            attachment_count += 1
                            filename = part.get_filename()
                            print(f"    ✅ 첨부파일 {attachment_count}: {filename}")
                    
                    print(f"[📊 첨부파일 개수] {attachment_count}개")
                    
                    if attachment_count > 0:
                        print(f"[📎 첨부파일 처리 시작] {attachment_count}개 발견")
                        attachments = extract_and_process_attachments_enhanced(msg, subject, email_id_str)
                        print(f"[✅ 첨부파일 처리 완료] {len(attachments)}개 처리됨")
                        
                        # 처리 결과 상세 로그
                        for i, att in enumerate(attachments):
                            print(f"  📎 {i+1}. {att.get('filename', 'Unknown')} ({att.get('type', 'unknown')})")
                            if att.get('document_summary'):
                                print(f"       📄 요약: {att['document_summary'][:50]}...")
                            if att.get('yolo_detections'):
                                print(f"       🤖 YOLO: {len(att['yolo_detections'])}개 객체")
                    else:
                        print(f"[ℹ️ 첨부파일 없음]")
                else:
                    print(f"[ℹ️ 단일 파트 메시지 - 첨부파일 없음]")
                    
            except Exception as e:
                print(f"[❗첨부파일 처리 오류] {str(e)}")
                import traceback
                traceback.print_exc()
                attachments = []
            # ===== 첨부파일 처리 추가 끝 =====

            # 분류 실행
            try:
                text_inputs = [body] + candidate_labels
                result = embed.text(text_inputs, model='nomic-embed-text-v1', task_type='classification')
                embedding_list = result['embeddings']
                email_embedding = [embedding_list[0]]
                label_embeddings = embedding_list[1:]
                scores = cosine_similarity(email_embedding, label_embeddings)[0]
                best_index = scores.argmax()
                classification_tag = candidate_labels[best_index]
                confidence = scores[best_index]
                print(f"[🏷️ 분류] {classification_tag} (신뢰도: {confidence:.3f})")
            except Exception as e:
                print("[⚠️ 분류 실패]", str(e))
                classification_tag = "unknown"

            # 요약 실행
            try:
                if not body:
                    summary_text = "(본문 없음)"
                else:
                    safe_text = body[:1000]
                    if len(safe_text) < 50:
                        summary_text = safe_text
                    else:
                        summary_text = summarizer(
                            safe_text,
                            max_length=80,
                            min_length=30,
                            do_sample=False
                        )[0]["summary_text"]
            except Exception as e:
                print("[⚠️ 요약 실패]", str(e))
                summary_text = body[:150] + "..." if body else "(요약 실패)"

            # 태그 추정
            typ, flag_data = mail.fetch(msg_id, "(FLAGS)")
            if flag_data and flag_data[0]:
                flags_bytes = flag_data[0]
                flags_str = flags_bytes.decode() if isinstance(flags_bytes, bytes) else str(flags_bytes)
            else:
                flags_str = ""

            tag = "받은"
            if "\\Important" in flags_str:
                tag = "중요"
            elif "\\Junk" in flags_str or "\\Spam" in flags_str:
                tag = "스팸"

            # ✅ 메일 객체에 첨부파일 정보 추가
            emails.append({
                "id": int(msg_id.decode()) if isinstance(msg_id, bytes) else int(msg_id),
                "subject": subject,
                "from": from_field,
                "date": date_str,
                "body": body[:1000],
                "tag": tag,
                "summary": summary_text,
                "classification": classification_tag,
                "attachments": attachments,  # ✅ 첨부파일 배열
                "has_attachments": len(attachments) > 0,  # ✅ 첨부파일 유무
                "attachment_summary": generate_enhanced_attachment_summary(attachments) if attachments else ""  # ✅ 첨부파일 요약
            })
            
            # 처리 완료 로그
            print(f"[✅ 메일 완료] {subject[:30]}... (첨부파일: {len(attachments)}개)")
            processed_count += 1

        # 백엔드에서도 날짜순 정렬 (최신 먼저)
        emails.sort(key=lambda x: x['date'], reverse=True)
        
        # 사용자별 세션에 메일 데이터 저장
        if user_key not in user_sessions:
            user_sessions[user_key] = {}
        
        user_sessions[user_key]['last_emails'] = emails
        user_sessions[user_key]['last_update'] = datetime.now().isoformat()
        
        print(f"[📊 결과] 사용자: {username}, 총 {processed_count}개 메일 처리 완료")
        if emails:
            print(f"[📅 범위] {emails[-1]['date']} ~ {emails[0]['date']}")

        return jsonify({
            "emails": emails,
            "user_session": user_key[:8] + "...",  # 디버그용
            "cache_info": f"세션에 {len(emails)}개 메일 저장됨"
        })

    except Exception as e:
        print("[❗에러 발생]", str(e))
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route('/api/chatbot', methods=['POST'])
def chatbot():
    try:
        data = request.get_json()
        user_input = data.get("user_input", "").strip()
        user_email = data.get("email", "")
        app_password = data.get("app_password", "")
        
        print(f"[🤖 챗봇 요청] 사용자: {user_email}, 입력: {user_input}")
        
        if not user_input:
            return jsonify({"error": "입력이 비어있습니다."}), 400
        
        # 사용자 세션 확인
        user_key = get_user_key(user_email)
        if user_key not in user_sessions:
            print(f"[⚠️ 세션 없음] {user_email} 사용자의 세션이 없습니다.")
            return jsonify({"error": "로그인이 필요합니다."}), 401
        
        # ✅ 1. 영어 Embedding 기반 분류 (기존 방식)
        candidate_labels = [
            "correct the vocabulary, spelling",
            "image generation using text", 
            "find something",
            "email search for a person"
        ]
        
        text_inputs = [user_input] + candidate_labels
        result = embed.text(text_inputs, model='nomic-embed-text-v1', task_type='classification')
        
        embedding_list = result['embeddings']
        email_embedding = [embedding_list[0]]
        label_embeddings = embedding_list[1:]
        
        scores = cosine_similarity(email_embedding, label_embeddings)[0]
        best_index = scores.argmax()
        embedding_score = scores[best_index]
        embedding_label = candidate_labels[best_index]
        
        # ✅ 2. 한국어 키워드 기반 분류
        user_input_lower = user_input.lower()
        
        korean_patterns = {
            "grammar": {
                "keywords": ["교정", "맞춤법", "문법", "틀렸", "고쳐", "수정"],
                "action": "grammar_correction"
            },
            "image": {
                "keywords": ["이미지", "그림", "사진", "그려", "만들어", "생성"],
                "action": "image_generation"
            },
            "person_search": {
                "keywords": ["님", "씨"],
                "required": ["메일", "이메일"],  # 둘 다 있어야 함
                "action": "person_search"
            },
            "general_search": {
                "keywords": ["찾아", "검색", "찾기"],
                "action": "email_search"
            }
        }
        
        korean_result = {"action": None, "confidence": 0.0, "matched_keywords": []}
        
        for pattern_name, pattern_info in korean_patterns.items():
            matched_keywords = []
            
            # 일반 키워드 매칭
            for keyword in pattern_info["keywords"]:
                if keyword in user_input_lower:
                    matched_keywords.append(keyword)
            
            # 필수 키워드 확인 (person_search용)
            if "required" in pattern_info:
                required_found = any(req in user_input_lower for req in pattern_info["required"])
                if not required_found:
                    continue
            
            # 신뢰도 계산 (매칭된 키워드 비율)
            if matched_keywords:
                confidence = len(matched_keywords) / len(pattern_info["keywords"])
                
                # person_search는 특별 처리 (필수 키워드 보너스)
                if pattern_name == "person_search" and "required" in pattern_info:
                    confidence += 0.3  # 보너스
                
                if confidence > korean_result["confidence"]:
                    korean_result = {
                        "action": pattern_info["action"],
                        "confidence": confidence,
                        "matched_keywords": matched_keywords
                    }
        
        print(f"[🔤 한국어 분석] {korean_result['action']} (신뢰도: {korean_result['confidence']:.3f})")
        print(f"[🌐 영어 분석] {embedding_label} (신뢰도: {embedding_score:.3f})")
        
        # ✅ 3. 최종 의도 결정 (더 높은 신뢰도 선택)
        
        # 영어 embedding 결과를 action으로 변환
        embedding_action_map = {
            "correct the vocabulary, spelling": "grammar_correction",
            "image generation using text": "image_generation", 
            "find something": "email_search",
            "email search for a person": "person_search"
        }
        
        embedding_action = embedding_action_map.get(embedding_label, "unknown")
        embedding_threshold = 0.25  # 임계값 낮춤
        
        # 최종 결정
        if korean_result["confidence"] >= 0.3 and korean_result["confidence"] > embedding_score:
            # 한국어 키워드 우선
            final_action = korean_result["action"]
            final_confidence = korean_result["confidence"]
            detection_method = "korean_keywords"
            
        elif embedding_score >= embedding_threshold:
            # 영어 embedding 사용
            final_action = embedding_action
            final_confidence = embedding_score
            detection_method = "english_embedding"
            
        else:
            # 둘 다 낮으면 unknown
            final_action = "unknown"
            final_confidence = max(korean_result["confidence"], embedding_score)
            detection_method = "low_confidence"
        
        print(f"[🎯 최종 결정] {final_action} (방법: {detection_method}, 신뢰도: {final_confidence:.3f})")
        
        # ✅ 4. 각 기능별 실행
        if final_action == "grammar_correction":
            response = handle_grammar_correction(user_input)
            
        elif final_action == "image_generation":
            response = handle_image_generation(user_input)
            
        elif final_action == "email_search":
            response = handle_general_search(user_input, user_email, app_password)
            
        elif final_action == "person_search":
            response = handle_person_search(user_input, user_email, app_password)
            
        else:
            response = """❓ 요청을 이해하지 못했습니다. 다른 표현을 시도해주세요.

🔧 **사용 가능한 기능들:**
• **문법/맞춤법 교정**: "이 문장 교정해주세요" / "correct this sentence"
• **이미지 생성**: "고양이 그림 그려줘" / "generate cat image"  
• **메일 검색**: "회의 관련 메일 찾아줘" / "find meeting emails"
• **사람별 메일**: "김철수님 메일 검색" / "search john@company.com emails"

💡 **Example / 예시:**
- 한국어: "안녕하세요. 제가 오늘 회의에 참석못할것 같습니다 교정해주세요"
- English: "correct the grammar: I can't attend meeting today"
- 혼합: "find 프로젝트 관련 emails" """
        
        return jsonify({
            "response": response,
            "action": final_action,
            "confidence": float(final_confidence),
            "detected_intent": final_action,
            "detection_method": detection_method,
            "debug_info": {
                "korean": f"{korean_result['action']} ({korean_result['confidence']:.3f})",
                "english": f"{embedding_action} ({embedding_score:.3f})",
                "final": f"{final_action} via {detection_method}"
            }
        }), 200
        
    except Exception as e:
        print("[❗챗봇 오류]", str(e))
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500
        
# ✅ 4. 각 기능별 핸들러 함수들
# handle_grammar_correction 함수를 이것으로 교체하세요

def handle_grammar_correction(user_input):
    """실제 작동하는 문법 및 맞춤법 교정 기능"""
    try:
        # 교정할 텍스트 추출
        correction_text = user_input
        
        # 불필요한 단어들 제거
        remove_words = ["교정해주세요", "교정해줘", "맞춤법", "문법", "correct", "spelling", "check", "fix"]
        for word in remove_words:
            correction_text = correction_text.replace(word, "").strip()
        
        if not correction_text:
            return "📝 **문법 및 맞춤법 교정**\n\n교정하고 싶은 텍스트를 입력해주세요.\n\n예시: '안녕하세요. 제가 오늘 회의에 참석못할것 같습니다' 교정해주세요"
        
        # ✅ 실제 HuggingFace API 사용한 교정
        hf_token = os.getenv("HF_TOKEN")
        if not hf_token:
            return f"📝 **문법 교정 결과**\n\n원본: {correction_text}\n\n⚠️ HF_TOKEN이 설정되지 않아 교정 서비스를 사용할 수 없습니다."
        
        try:
            from huggingface_hub import InferenceClient
            
            client = InferenceClient(
                model="Qwen/Qwen2.5-7B-Instruct",
                token=hf_token
            )
            
            # 더 나은 프롬프트 작성
            prompt = f"""다음 텍스트의 맞춤법, 문법, 띄어쓰기를 교정해주세요. 자연스럽고 정확한 한국어/영어로 수정해주세요.

원본 텍스트:
"{correction_text}"

교정 지침:
1. 맞춤법 오류 수정
2. 문법 오류 수정  
3. 띄어쓰기 수정
4. 자연스러운 표현으로 개선
5. 원래 의미는 유지

교정된 텍스트:"""
            
            messages = [
                {"role": "system", "content": "당신은 전문 교정 편집자입니다. 주어진 텍스트의 맞춤법, 문법, 띄어쓰기를 정확하게 교정합니다."},
                {"role": "user", "content": prompt}
            ]
            
            response = client.chat_completion(
                messages=messages,
                max_tokens=300,
                temperature=0.3
            )
            
            corrected_text = response.choices[0].message.content.strip()
            
            # 교정 결과 분석
            changes_made = []
            
            # 간단한 변화 감지
            if len(corrected_text) != len(correction_text):
                changes_made.append("길이 변경")
            if corrected_text != correction_text:
                changes_made.append("내용 수정")
            
            print(f"[✅ 실제 문법 교정 완료] 원본: {len(correction_text)}자 -> 교정: {len(corrected_text)}자")
            
            return f"""📝 **문법 및 맞춤법 교정 완료**

**원본:**
{correction_text}

**교정된 텍스트:**
{corrected_text}

**변경사항:**
{', '.join(changes_made) if changes_made else '수정할 내용이 없습니다.'}

✅ **AI 교정이 완료되었습니다!**
💡 교정 결과를 검토한 후 사용하시기 바랍니다."""
            
        except Exception as e:
            print(f"[❗교정 API 오류] {str(e)}")
            
            # ✅ 간단한 규칙 기반 교정으로 fallback
            simple_corrections = {
                # 자주 틀리는 맞춤법
                "데이타": "데이터",
                "컴퓨타": "컴퓨터", 
                "셋팅": "설정",
                "미팅": "회의",
                "어플리케이션": "애플리케이션",
                "어플": "앱",
                
                # 띄어쓰기
                "안녕하세요.": "안녕하세요. ",
                "입니다.": "입니다. ",
                "합니다.": "합니다. ",
                
                # 자주 틀리는 표현
                "해야되는": "해야 하는",
                "할수있는": "할 수 있는",
                "못할것": "못할 것",
                "참석못할": "참석하지 못할"
            }
            
            corrected_simple = correction_text
            applied_corrections = []
            
            for wrong, correct in simple_corrections.items():
                if wrong in corrected_simple:
                    corrected_simple = corrected_simple.replace(wrong, correct)
                    applied_corrections.append(f"'{wrong}' → '{correct}'")
            
            if applied_corrections:
                return f"""📝 **간단 맞춤법 교정**

**원본:**
{correction_text}

**교정된 텍스트:**
{corrected_simple}

**적용된 교정:**
{chr(10).join('• ' + correction for correction in applied_corrections)}

⚠️ **참고:** AI 교정 서비스 연결 실패로 기본 규칙만 적용했습니다."""
            else:
                return f"""📝 **교정 검토 완료**

**입력된 텍스트:**
{correction_text}

✅ **현재 텍스트에서 명백한 오류를 발견하지 못했습니다.**

💡 **참고:**
• AI 교정 서비스 연결에 실패했습니다
• 더 정확한 교정을 위해 다시 시도해보세요
• 특정 부분이 의심스럽다면 구체적으로 지적해주세요"""
            
    except Exception as e:
        print(f"[❗문법 교정 오류] {str(e)}")
        return "❌ 문법 교정 처리 중 오류가 발생했습니다."
# handle_image_generation 함수를 이것으로 교체하세요

def handle_image_generation(user_input):
    """실제 작동하는 이미지 생성 기능 (HuggingFace API 사용)"""
    try:
        # 이미지 생성 프롬프트 추출
        image_prompt = user_input
        
        # 불필요한 단어들 제거
        remove_words = ["이미지 생성해주세요", "이미지 생성", "그려줘", "그림", "image generation", "generate", "만들어"]
        for word in remove_words:
            image_prompt = image_prompt.replace(word, "").strip()
        
        if not image_prompt:
            return "🎨 **이미지 생성**\n\n생성하고 싶은 이미지에 대한 설명을 입력해주세요.\n\n예시:\n• '아름다운 석양과 바다'\n• '귀여운 고양이가 놀고 있는 모습'\n• 'A beautiful sunset over the ocean'"
        
        # ✅ 실제 HuggingFace 이미지 생성 API 사용
        hf_token = os.getenv("HF_TOKEN")
        if not hf_token:
            return f"🎨 **이미지 생성**\n\n요청된 이미지: '{image_prompt}'\n\n⚠️ HF_TOKEN이 설정되지 않아 이미지 생성을 할 수 없습니다."
        
        try:
            from huggingface_hub import InferenceClient
            import base64
            import time
            
            # ✅ Stable Diffusion 모델 사용
            client = InferenceClient(
                model="runwayml/stable-diffusion-v1-5",
                token=hf_token
            )
            
            # 한국어를 영어로 번역 (더 좋은 결과를 위해)
            korean_to_english = {
                "고양이": "cute cat",
                "강아지": "cute dog", 
                "꽃": "beautiful flowers",
                "바다": "ocean and waves",
                "산": "mountains and nature",
                "석양": "beautiful sunset",
                "하늘": "blue sky with clouds",
                "숲": "forest and trees",
                "도시": "modern city",
                "자동차": "modern car",
                "집": "beautiful house",
                "사람": "person",
                "음식": "delicious food",
                "케이크": "beautiful cake"
            }
            
            # 영어 프롬프트 생성
            english_prompt = image_prompt
            if any(ord(char) > 127 for char in image_prompt):  # 한국어 포함 여부 확인
                for korean, english in korean_to_english.items():
                    if korean in image_prompt:
                        english_prompt = english_prompt.replace(korean, english)
                
                # 매칭되지 않은 한국어가 있으면 기본 프롬프트 생성
                if any(ord(char) > 127 for char in english_prompt):
                    english_prompt = f"a beautiful {image_prompt}"
            
            # 프롬프트 품질 개선
            enhanced_prompt = f"{english_prompt}, high quality, detailed, beautiful, artistic"
            
            print(f"[🎨 이미지 생성 시작] '{image_prompt}' -> '{enhanced_prompt}'")
            
            # ✅ 실제 이미지 생성
            image_bytes = client.text_to_image(
                prompt=enhanced_prompt,
                height=512,
                width=512,
                num_inference_steps=20
            )
            
            # 이미지를 base64로 인코딩 (웹에서 표시하기 위해)
            image_base64 = base64.b64encode(image_bytes).decode()
            
            # 이미지 파일로 저장
            timestamp = int(time.time())
            filename = f"generated_image_{timestamp}.png"
            filepath = os.path.join(ATTACHMENT_FOLDER, filename)
            
            # 디렉토리 확인 및 생성
            os.makedirs(ATTACHMENT_FOLDER, exist_ok=True)
            
            with open(filepath, "wb") as f:
                f.write(image_bytes)
            
            print(f"[✅ 이미지 생성 완료] 파일 저장: {filepath}")
            
            return f"""🎨 **이미지 생성 완료!**

📝 **요청:** '{image_prompt}'
🖼️ **생성된 이미지:** {filename}
📁 **저장 위치:** /static/attachments/{filename}
🌐 **웹 주소:** http://localhost:5001/static/attachments/{filename}

✅ **성공!** 이미지가 생성되어 저장되었습니다.

🔗 **이미지 정보:**
- 파일명: {filename}
- 크기: 512x512 픽셀
- 프롬프트: "{enhanced_prompt}"
- 저장 시간: {time.strftime('%Y-%m-%d %H:%M:%S')}

💡 **사용 방법:**
1. 위 웹 주소를 브라우저에서 열어 이미지 확인
2. 파일 탐색기에서 static/attachments 폴더 확인
3. 더 구체적인 설명으로 다른 이미지 생성 가능

🎯 **팁:** 더 구체적인 설명을 하면 더 좋은 이미지를 얻을 수 있습니다!"""

        except Exception as e:
            error_msg = str(e)
            print(f"[❗이미지 생성 실패] {error_msg}")
            
            # ✅ 구체적인 오류 대응
            if "rate limit" in error_msg.lower():
                return f"""🎨 **이미지 생성 - 일시적 제한**

요청된 이미지: '{image_prompt}'

⏳ **잠시 후 다시 시도해주세요**
오류: API 요청 한도 초과

💡 **대안:**
• 1-2분 후 다시 시도
• 다른 키워드로 시도
• 영어로 입력: "{image_prompt} in English" """
                
            elif "unauthorized" in error_msg.lower() or "token" in error_msg.lower():
                return f"""🎨 **이미지 생성 - 인증 오류**

요청된 이미지: '{image_prompt}'

🔑 **인증 문제 발생**
오류: HuggingFace 토큰 인증 실패

💡 **해결방법:**
• 관리자에게 HF_TOKEN 확인 요청
• 토큰이 올바른지 확인"""
                
            else:
                return f"""🎨 **이미지 생성 실패**

요청된 이미지: '{image_prompt}'

❌ **생성 실패**
오류: {error_msg}

💡 **다른 방법 시도:**
• 더 간단한 설명 사용: "cat", "flower", "sunset"
• 영어로 입력해보세요
• 특수문자 제거 후 재시도

📝 **예시:**
- "beautiful landscape" 
- "cute animal"
- "modern building" """
            
    except Exception as e:
        print(f"[❗이미지 생성 핸들러 오류] {str(e)}")
        return "❌ 이미지 생성 처리 중 오류가 발생했습니다."
    
def handle_general_search(user_input, user_email, app_password):
    """일반 키워드 메일 검색 (개선된 버전)"""
    try:
        print(f"[🔍 일반 검색 시작] 입력: '{user_input}', 사용자: {user_email}")
        
        # 검색 키워드 추출 개선
        search_keywords = user_input.lower()
        
        # 불필요한 단어들 제거
        remove_words = ["찾아줘", "찾아주세요", "검색해줘", "검색", "find", "search", "메일", "이메일", "email"]
        for word in remove_words:
            search_keywords = search_keywords.replace(word, "").strip()
        
        print(f"[🎯 추출된 키워드] '{search_keywords}'")
        
        if not search_keywords:
            return "🔍 **메일 검색**\n\n검색하고 싶은 키워드를 입력해주세요.\n\n예시:\n• '회의 관련 메일 찾아줘'\n• '프로젝트 업데이트 검색'\n• '급한 메일 찾기'"
        
        # 실제 메일 검색 로직
        try:
            # 메일 서버 연결
            print("[📧 메일 서버 연결 시작]")
            mail = imaplib.IMAP4_SSL("imap.gmail.com")
            mail.login(user_email, app_password)
            mail.select("inbox")
            print("[✅ 메일 서버 연결 성공]")
            
            # 더 많은 메일 검색 (범위 확대)
            N = 50  # 50개로 증가
            status, data_result = mail.search(None, "ALL")
            all_mail_ids = data_result[0].split()
            mail_ids = all_mail_ids[-N:]
            
            print(f"[📊 검색 범위] 총 {len(all_mail_ids)}개 중 최근 {len(mail_ids)}개 검색")
            
            found_emails = []
            processed_count = 0
            
            for msg_id in mail_ids:
                try:
                    _, msg_data = mail.fetch(msg_id, "(RFC822)")
                    if not msg_data or not msg_data[0]:
                        continue
                    
                    msg = email_module.message_from_bytes(msg_data[0][1])
                    processed_count += 1
                    
                    # 제목 디코딩 (기존 summary 함수와 같은 방식)
                    raw_subject = msg.get("Subject", "")
                    try:
                        decoded_parts = decode_header(raw_subject)
                        if decoded_parts and decoded_parts[0]:
                            decoded_subject = decoded_parts[0]
                            subject_bytes = decoded_subject[0]
                            subject_encoding = decoded_subject[1]
                            
                            if isinstance(subject_bytes, bytes):
                                if subject_encoding is None:
                                    subject_encoding = 'utf-8'
                                try:
                                    subject = subject_bytes.decode(subject_encoding)
                                except (UnicodeDecodeError, LookupError):
                                    subject = subject_bytes.decode('utf-8', errors='ignore')
                            else:
                                subject = str(subject_bytes)
                        else:
                            subject = "(제목 없음)"
                    except Exception as e:
                        subject = raw_subject if raw_subject else "(제목 없음)"
                    
                    # 발신자 정보
                    name, addr = parseaddr(msg.get("From"))
                    from_field = f"{name} <{addr}>" if name else addr
                    
                    # 날짜 정보
                    raw_date = msg.get("Date", "")
                    try:
                        date_obj = parsedate_to_datetime(raw_date)
                        date_str = date_obj.strftime("%Y-%m-%d %H:%M")
                    except:
                        date_str = raw_date[:16] if len(raw_date) >= 16 else raw_date
                    
                    # 본문 추출 (간단한 버전)
                    body = ""
                    try:
                        if msg.is_multipart():
                            for part in msg.walk():
                                if part.get_content_type() == "text/plain" and not part.get("Content-Disposition"):
                                    charset = part.get_content_charset() or "utf-8"
                                    body += part.get_payload(decode=True).decode(charset, errors="ignore")
                        else:
                            charset = msg.get_content_charset() or "utf-8"
                            body = msg.get_payload(decode=True).decode(charset, errors="ignore")
                        body = body.strip()[:200]  # 처음 200자만
                    except Exception as e:
                        body = ""
                    
                    # 개선된 키워드 검색 (제목, 발신자, 본문에서 모두 검색)
                    search_in = f"{subject} {from_field} {body}".lower()
                    
                    # 여러 키워드 중 하나라도 매칭되면 포함
                    keywords = search_keywords.split()
                    if any(keyword in search_in for keyword in keywords):
                        found_emails.append({
                            "subject": subject[:60] + "..." if len(subject) > 60 else subject,
                            "from": from_field[:40] + "..." if len(from_field) > 40 else from_field,
                            "date": date_str,
                            "preview": body[:100] + "..." if len(body) > 100 else body
                        })
                        
                        print(f"[✅ 매칭] {subject[:30]}...")
                        
                        if len(found_emails) >= 8:  # 최대 8개까지
                            break
                            
                except Exception as e:
                    print(f"[⚠️ 메일 처리 오류] {str(e)}")
                    continue
            
            mail.close()
            mail.logout()
            
            print(f"[📊 검색 완료] {processed_count}개 처리, {len(found_emails)}개 발견")
            
            if found_emails:
                result = f"🔍 **검색 결과**\n\n키워드: '{search_keywords}'\n검색된 메일: {len(found_emails)}개 (총 {processed_count}개 중)\n\n"
                for i, mail_info in enumerate(found_emails, 1):
                    result += f"**{i}. {mail_info['subject']}**\n"
                    result += f"📤 {mail_info['from']}\n"
                    result += f"📅 {mail_info['date']}\n"
                    if mail_info['preview']:
                        result += f"💬 {mail_info['preview']}\n"
                    result += "\n"
                result += "💡 더 정확한 검색을 위해 구체적인 키워드를 사용해보세요."
                return result
            else:
                return f"🔍 **검색 결과**\n\n키워드: '{search_keywords}'\n검색 범위: 최근 {processed_count}개 메일\n\n❌ 관련된 메일을 찾을 수 없습니다.\n\n💡 **검색 팁:**\n• 다른 키워드로 시도해보세요\n• 발신자 이름이나 회사명 사용\n• 메일 제목의 핵심 단어 사용\n• 영어/한국어 모두 시도"
                
        except Exception as e:
            print(f"[❗메일 검색 오류] {str(e)}")
            return f"❌ 메일 검색 중 오류가 발생했습니다.\n\n오류 내용: {str(e)}\n\n💡 로그인 정보나 네트워크 연결을 확인해주세요."
        
    except Exception as e:
        print(f"[❗일반 검색 오류] {str(e)}")
        return "❌ 검색 처리 중 오류가 발생했습니다."

def handle_person_search(user_input, user_email, app_password):
    """특정 사람 메일 검색 (개선된 버전)"""
    try:
        print(f"[👤 사람 검색 시작] 입력: '{user_input}'")
        
        # Qwen을 이용해 사람 이름/이메일 추출
        search_target = extract_search_target_with_qwen(user_input)
        print(f"[🎯 추출된 대상] '{search_target}'")
        
        # Qwen 실패 시 간단한 추출 방법
        if not search_target or len(search_target.strip()) < 2:
            # 간단한 이름/이메일 추출
            words = user_input.split()
            potential_targets = []
            
            for word in words:
                # 이메일 주소 패턴
                if "@" in word and "." in word:
                    potential_targets.append(word)
                # 한국어 이름 패턴 (2-4글자)
                elif len(word) >= 2 and len(word) <= 4 and word.replace(" ", "").isalpha():
                    potential_targets.append(word)
            
            if potential_targets:
                search_target = potential_targets[0]
            else:
                return "👤 **사람별 메일 검색**\n\n찾고 싶은 사람의 이름이나 이메일 주소를 명확히 알려주세요.\n\n예시:\n• '김철수님의 메일'\n• 'john@company.com 메일'\n• '홍길동 교수님 메일'"
        
        print(f"[🔍 최종 검색 대상] '{search_target}'")
        
        try:
            # 메일 서버 연결
            mail = imaplib.IMAP4_SSL("imap.gmail.com")
            mail.login(user_email, app_password)
            mail.select("inbox")
            
            # 더 많은 메일 검색
            N = 100  # 100개로 증가
            status, data_result = mail.search(None, "ALL")
            all_mail_ids = data_result[0].split()
            mail_ids = all_mail_ids[-N:]
            
            print(f"[📊 검색 범위] 최근 {len(mail_ids)}개 메일에서 검색")
            
            found_emails = []
            processed_count = 0
            
            for msg_id in mail_ids:
                try:
                    _, msg_data = mail.fetch(msg_id, "(RFC822)")
                    if not msg_data or not msg_data[0]:
                        continue
                    
                    msg = email_module.message_from_bytes(msg_data[0][1])
                    processed_count += 1
                    
                    # 발신자 정보 추출
                    from_header = msg.get("From", "")
                    name, addr = parseaddr(from_header)
                    from_field = f"{name} <{addr}>" if name else addr
                    
                    # 제목 추출 (간단한 방법)
                    subject = str(msg.get("Subject", ""))[:80]
                    
                    # 날짜 추출
                    date_field = str(msg.get("Date", ""))[:25]
                    
                    # 검색 대상이 발신자 정보에 포함되는지 확인 (대소문자 무시, 부분 매칭)
                    search_lower = search_target.lower()
                    from_lower = from_field.lower()
                    
                    # 더 관대한 매칭
                    if (search_lower in from_lower or 
                        any(part.strip() in from_lower for part in search_lower.split() if part.strip()) or
                        (len(search_lower) >= 3 and search_lower in from_lower.replace(" ", ""))):
                        
                        found_emails.append({
                            "subject": subject,
                            "from": from_field,
                            "date": date_field
                        })
                        
                        print(f"[✅ 매칭] {from_field} -> {subject[:30]}...")
                        
                        if len(found_emails) >= 10:  # 최대 10개까지
                            break
                            
                except Exception as e:
                    continue
            
            mail.close()
            mail.logout()
            
            print(f"[📊 사람 검색 완료] {processed_count}개 처리, {len(found_emails)}개 발견")
            
            if found_emails:
                result = f"👤 **사람별 메일 검색 결과**\n\n검색 대상: '{search_target}'\n발견된 메일: {len(found_emails)}개 (총 {processed_count}개 중)\n\n"
                for i, mail_info in enumerate(found_emails, 1):
                    result += f"**{i}. {mail_info['subject']}**\n"
                    result += f"📤 {mail_info['from']}\n"
                    result += f"📅 {mail_info['date']}\n\n"
                result += "💡 특정 메일을 자세히 보려면 메일 리스트에서 확인하세요."
                return result
            else:
                return f"👤 **사람별 메일 검색 결과**\n\n검색 대상: '{search_target}'\n검색 범위: 최근 {processed_count}개 메일\n\n❌ 해당 사람의 메일을 찾을 수 없습니다.\n\n💡 **검색 팁:**\n• 정확한 이름 사용: '{search_target}' → 다른 표기법 시도\n• 이메일 주소로 시도\n• 성이나 이름만으로 시도\n• 영문/한글 이름 모두 시도"
                
        except Exception as e:
            print(f"[❗사람 검색 오류] {str(e)}")
            return f"❌ 사람별 메일 검색 중 오류가 발생했습니다.\n\n오류: {str(e)}"
        
    except Exception as e:
        print(f"[❗사람 검색 핸들러 오류] {str(e)}")
        return "❌ 사람 검색 처리 중 오류가 발생했습니다."

@app.route('/api/test', methods=['POST'])
def test():
    data = request.get_json()
    text = data.get("text", "")
    email = data.get("email", "")
    
    user_key = get_user_key(email) if email else "anonymous"
    
    return jsonify({
        "message": f"✅ 백엔드 정상 작동: {text[:20]}...",
        "user_session": user_key[:8] + "..." if email else "no_session"
    })

@app.route("/api/send", methods=["POST"])
def send_email():
    try:
        data = request.get_json()
        print("✅ 받은 데이터:", data)

        sender_email = data["email"]
        app_password = data["app_password"]
        to = data["to"]
        subject = data["subject"]
        body = data["body"]

        # 사용자 세션 확인
        user_key = get_user_key(sender_email)
        if user_key not in user_sessions:
            return jsonify({"error": "로그인이 필요합니다."}), 401

        msg = MIMEText(body)
        msg["Subject"] = subject
        msg["From"] = sender_email
        msg["To"] = to

        server = smtplib.SMTP_SSL("smtp.gmail.com", 465)
        server.login(sender_email, app_password)
        server.send_message(msg)
        server.quit()

        print(f"[📤 메일 전송 성공] 사용자: {sender_email}, 수신자: {to}")

        return jsonify({"message": "✅ 메일 전송 성공"}), 200

    except Exception as e:
        print("[❗메일 전송 실패]", str(e))
        return jsonify({"error": str(e)}), 500

@app.route('/api/session-info', methods=['GET'])
def session_info():
    """현재 활성 세션 정보 반환 (디버그용)"""
    return jsonify({
        "active_sessions": len(user_sessions),
        "session_keys": [key[:8] + "..." for key in user_sessions.keys()],
        "yolo_model_loaded": yolo_model is not None
    })

@app.route('/', methods=['GET'])
def health_check():
    return "✅ 백엔드 정상 작동 중 (사용자 세션 분리 적용)\n{yolo_status}", 200

# ✅ 새로운 첨부파일 정보 API 추가
@app.route('/api/attachment-info', methods=['POST'])
def attachment_info():
    """특정 메일의 첨부파일 상세 정보 반환"""
    try:
        data = request.get_json()
        email_id = data.get("email_id")
        user_email = data.get("email", "")
        
        # 사용자 세션 확인
        user_key = get_user_key(user_email)
        if user_key not in user_sessions:
            return jsonify({"error": "로그인이 필요합니다."}), 401
        
        # 세션에서 해당 메일 찾기
        last_emails = user_sessions[user_key].get('last_emails', [])
        target_email = None
        
        for email_data in last_emails:
            if email_data.get('id') == email_id:
                target_email = email_data
                break
        
        if not target_email:
            return jsonify({"error": "해당 메일을 찾을 수 없습니다."}), 404
        
        attachments = target_email.get('attachments', [])
        
        return jsonify({
            "success": True,
            "email_id": email_id,
            "subject": target_email.get('subject', ''),
            "attachments": attachments,
            "attachment_count": len(attachments),
            "has_yolo_detections": any(att.get('yolo_detections') for att in attachments)
        })
        
    except Exception as e:
        print(f"[❗첨부파일 정보 오류] {str(e)}")
        return jsonify({"error": str(e)}), 500


# to 대쉬보드용 추가 코드

# 할일 추출을 위한 키워드 패턴들
TODO_KEYWORDS = {
    'meeting': ['회의', '미팅', 'meeting', '컨퍼런스', '세미나', '면담', '상담'],
    'deadline': ['마감', '제출', '완료', '끝내', 'deadline', 'due', '기한', '까지'],
    'task': ['작업', '업무', '처리', '진행', '해야', '할것', 'task', 'work', 'todo'],
    'event': ['행사', '이벤트', 'event', '파티', '모임', '약속', '일정'],
    'reminder': ['알림', 'reminder', '잊지말', '기억', '체크', '확인']
}

# 날짜/시간 추출 패턴
DATE_PATTERNS = [
    r'(\d{4})년\s*(\d{1,2})월\s*(\d{1,2})일',  # 2024년 12월 25일
    r'(\d{1,2})월\s*(\d{1,2})일',              # 12월 25일
    r'(\d{1,2})/(\d{1,2})',                    # 12/25
    r'(\d{4}-\d{1,2}-\d{1,2})',               # 2024-12-25
    r'(오늘|내일|모레)',                        # 상대적 날짜
    r'(다음주|이번주|다다음주)',                 # 상대적 주
    r'(월요일|화요일|수요일|목요일|금요일|토요일|일요일)'  # 요일
]

TIME_PATTERNS = [
    r'(\d{1,2}):(\d{2})',                      # 14:30
    r'(\d{1,2})시\s*(\d{1,2})?분?',           # 2시 30분
    r'(오전|오후)\s*(\d{1,2})시',              # 오전 10시
]

def extract_todos_from_email_improved(email_body, email_subject, email_from, email_date):
    """개선된 이메일 할일 추출 함수 - 고유 ID 추가"""
    try:
        print(f"[📋 개선된 할일 추출] {email_subject[:30]}...")
        
        full_text = f"{email_subject} {email_body}"
        todos = []
        
        # ✅ 고유 ID 생성을 위한 기준 시간
        import time
        base_timestamp = int(time.time() * 1000)
        
        # 1. 회의/미팅 스마트 추출
        try:
            meeting_todos = extract_meetings_improved(full_text, email_from, email_date, email_subject, base_timestamp)
            todos.extend(meeting_todos)
        except Exception as e:
            print(f"[⚠️ 회의 추출 오류] {str(e)}")

        # 2. 마감일/데드라인 스마트 추출  
        deadline_todos = extract_deadlines_improved(full_text, email_from, email_date, email_subject, base_timestamp + 100)
        todos.extend(deadline_todos)
        
        # 3. 일반 할일 스마트 추출
        task_todos = extract_general_tasks_improved(full_text, email_from, email_date, email_subject, base_timestamp + 200)
        todos.extend(task_todos)
        
        # 4. 이벤트/행사 스마트 추출
        event_todos = extract_events_improved(full_text, email_from, email_date, email_subject, base_timestamp + 300)
        todos.extend(event_todos)
        
        # 중복 제거 및 우선순위 설정
        todos = deduplicate_todos_improved(todos)
        todos = assign_priority(todos)
        
        print(f"[✅ 개선된 할일 추출 완료] {len(todos)}개 발견")
        
        return {
            'success': True,
            'todos': todos,
            'total_count': len(todos),
            'extraction_method': 'improved_ai_analysis'
        }
        
    except Exception as e:
        print(f"[❗할일 추출 오류] {str(e)}")
        return {
            'success': False,
            'todos': [],
            'error': str(e)
        }
    
# 이 함수를 extract_deadlines_improved 함수 바로 위에 추가하세요
def extract_meetings_improved(text, sender, email_date, email_subject, base_id):
    """개선된 회의/미팅 추출 - 고유 ID 추가"""
    meetings = []
    
    meeting_keywords = ['회의', '미팅', 'meeting', '면담', '상담', '컨퍼런스', '세미나']
    
    for keyword in meeting_keywords:
        if keyword.lower() in text.lower():
            meeting_title = generate_smart_title(text, keyword, email_subject, 'meeting')
            meeting_date = extract_smart_date(text) or '2024-12-27'
            meeting_time = extract_smart_time(text) or '14:00'
            
            meeting = {
                'id': base_id + len(meetings),  # ✅ 고유 ID
                'type': 'meeting',
                'title': meeting_title,
                'description': f"{sender}님과의 {keyword}",
                'date': meeting_date,
                'time': meeting_time,
                'priority': 'high',
                'status': 'pending',
                'editable_date': True,  # ✅ 날짜 편집 가능
                'source_email': {
                    'from': sender,
                    'subject': email_subject,
                    'date': email_date,
                    'type': 'meeting_invitation'
                }
            }
            meetings.append(meeting)
            print(f"[🤝 회의 추출] {meeting_title} (ID: {meeting['id']})")
            break
    
    return meetings

def extract_deadlines_improved(text, sender, email_date, email_subject, base_id):
    """개선된 마감일 추출 - 고유 ID 추가"""
    deadlines = []
    
    deadline_keywords = ['마감', '제출', '완료', 'deadline', 'due', '기한', '까지', 'submit']
    
    for keyword in deadline_keywords:
        if keyword.lower() in text.lower():
            deadline_title = generate_smart_title(text, keyword, email_subject, 'deadline')
            deadline_date = extract_smart_date(text) or '2024-12-28'
            
            deadline = {
                'id': base_id + len(deadlines),  # ✅ 고유 ID
                'type': 'deadline',
                'title': deadline_title,
                'description': f"{sender}님이 요청한 마감 업무",
                'date': deadline_date,
                'time': None,
                'priority': 'high',
                'status': 'pending',
                'editable_date': True,  # ✅ 날짜 편집 가능
                'source_email': {
                    'from': sender,
                    'subject': email_subject,
                    'date': email_date,
                    'type': 'deadline_notice'
                }
            }
            deadlines.append(deadline)
            print(f"[⏰ 마감일 추출] {deadline_title} (ID: {deadline['id']})")
            break
    
    return deadlines

def extract_general_tasks_improved(text, sender, email_date, email_subject, base_id):
    """개선된 일반 업무 추출 - 고유 ID 추가"""
    tasks = []
    
    task_patterns = [
        (r'([^\n\.]{10,80})\s*(해주세요|해주시기|부탁드립니다|요청드립니다)', 'korean_request'),
        (r'([^\n\.]{10,80})\s*(확인|검토|처리|진행)\s*(해주세요|부탁|필요)', 'korean_action'),
        (r'(please|kindly)\s+([^\n\.]{10,80})', 'english_request'),
        (r'([^\n\.]{10,80})\s+(please|kindly)', 'english_request_after'),
        (r'(could you|can you|would you)\s+([^\n\.]{10,80})', 'english_question'),
    ]
    
    for pattern, pattern_type in task_patterns:
        matches = re.finditer(pattern, text, re.IGNORECASE)
        for match in matches:
            if pattern_type == 'english_request':
                task_name = match.group(2).strip()
            elif pattern_type == 'english_request_after':
                task_name = match.group(1).strip()
            elif pattern_type == 'english_question':
                task_name = match.group(2).strip()
            else:
                task_name = match.group(1).strip()
            
            if len(task_name) > 10 and not is_meaningless_text(task_name):
                clean_title = clean_task_title(task_name)
                
                task = {
                    'id': base_id + len(tasks),  # ✅ 고유 ID
                    'type': 'task',
                    'title': clean_title,
                    'description': f"{sender}님이 요청한 업무",
                    'date': None,
                    'time': None,
                    'priority': 'medium',
                    'status': 'pending',
                    'editable_date': True,  # ✅ 날짜 편집 가능
                    'source_email': {
                        'from': sender,
                        'subject': email_subject,
                        'date': email_date,
                        'type': 'task_request'
                    }
                }
                tasks.append(task)
                print(f"[📋 업무 추출] {clean_title} (ID: {task['id']})")
                
                if len(tasks) >= 3:
                    break
    
    return tasks

def extract_events_improved(text, sender, email_date, email_subject, base_id):
    """개선된 이벤트 추출 - 고유 ID 추가"""
    events = []
    
    event_keywords = ['행사', '이벤트', 'event', '파티', '모임', '세미나', '워크샵', 'workshop']
    
    for keyword in event_keywords:
        if keyword.lower() in text.lower():
            event_title = generate_smart_title(text, keyword, email_subject, 'event')
            
            event = {
                'id': base_id + len(events),  # ✅ 고유 ID
                'type': 'event',
                'title': event_title,
                'description': f"{sender}님이 알린 {keyword}",
                'date': extract_smart_date(text) or '2024-12-29',
                'time': extract_smart_time(text) or '18:00',
                'priority': 'medium',
                'status': 'pending',
                'editable_date': True,  # ✅ 날짜 편집 가능
                'source_email': {
                    'from': sender,
                    'subject': email_subject,
                    'date': email_date,
                    'type': 'event_notification'
                }
            }
            events.append(event)
            print(f"[🎉 이벤트 추출] {event_title} (ID: {event['id']})")
            break
    
    return events

def generate_smart_title(text, keyword, email_subject, todo_type):
    """스마트한 제목 생성"""
    
    # 1. 이메일 제목에서 키워드 관련 부분 찾기
    subject_words = email_subject.split()
    
    # 2. 제목에 키워드가 있으면 제목 사용
    if keyword.lower() in email_subject.lower():
        return email_subject[:60]
    
    # 3. 본문에서 키워드 주변 문장 찾기
    sentences = text.split('.')
    for sentence in sentences:
        if keyword.lower() in sentence.lower() and len(sentence.strip()) > 10:
            clean_sentence = sentence.strip()
            if len(clean_sentence) > 60:
                clean_sentence = clean_sentence[:60] + "..."
            return clean_sentence
    
    # 4. 기본 제목 생성
    type_names = {
        'meeting': '회의',
        'deadline': '마감일',
        'task': '업무',
        'event': '이벤트'
    }
    
    base_name = type_names.get(todo_type, '할일')
    return f"{base_name}: {email_subject[:40]}"

def extract_smart_date(text):
    """스마트 날짜 추출"""
    
    # 한국어 날짜 패턴들
    korean_patterns = [
        r'(\d{4})년\s*(\d{1,2})월\s*(\d{1,2})일',
        r'(\d{1,2})월\s*(\d{1,2})일',
        r'(\d{1,2})/(\d{1,2})',
    ]
    
    for pattern in korean_patterns:
        match = re.search(pattern, text)
        if match:
            try:
                if len(match.groups()) == 3:  # 년월일
                    year, month, day = match.groups()
                    return f"{year}-{month.zfill(2)}-{day.zfill(2)}"
                elif len(match.groups()) == 2:  # 월일
                    month, day = match.groups()
                    year = datetime.now().year
                    return f"{year}-{month.zfill(2)}-{day.zfill(2)}"
            except:
                continue
    
    # 상대적 날짜
    today = datetime.now()
    if '오늘' in text:
        return today.strftime('%Y-%m-%d')
    elif '내일' in text:
        return (today + timedelta(days=1)).strftime('%Y-%m-%d')
    elif '다음주' in text:
        return (today + timedelta(days=7)).strftime('%Y-%m-%d')
    
    return None

def extract_smart_time(text):
    """스마트 시간 추출"""
    
    # 시간 패턴들
    time_patterns = [
        r'(\d{1,2}):(\d{2})',  # 14:30
        r'오전\s*(\d{1,2})시',  # 오전 10시
        r'오후\s*(\d{1,2})시',  # 오후 2시
        r'(\d{1,2})시\s*(\d{1,2})?분?',  # 2시 30분
    ]
    
    for pattern in time_patterns:
        match = re.search(pattern, text)
        if match:
            if '오전' in pattern:
                hour = int(match.group(1))
                return f"{hour:02d}:00"
            elif '오후' in pattern:
                hour = int(match.group(1))
                if hour != 12:
                    hour += 12
                return f"{hour:02d}:00"
            elif ':' in match.group(0):
                return match.group(0)
            else:
                hour = int(match.group(1))
                minute = match.group(2) if len(match.groups()) > 1 and match.group(2) else "00"
                return f"{hour:02d}:{minute.zfill(2)}"
    
    return None

def is_meaningless_text(text):
    """의미없는 텍스트인지 확인"""
    meaningless_patterns = [
        r'^[^a-zA-Z가-힣]*$',  # 문자가 없음
        r'^(please|kindly|확인|검토|처리)$',  # 단일 키워드만
        r'^.{1,5}$',  # 너무 짧음
    ]
    
    for pattern in meaningless_patterns:
        if re.match(pattern, text.strip(), re.IGNORECASE):
            return True
    
    return False

def clean_task_title(title):
    """할일 제목 정리"""
    
    # 불필요한 단어들 제거
    remove_words = ['해주세요', '부탁드립니다', '요청드립니다', 'please', 'kindly']
    
    clean_title = title
    for word in remove_words:
        clean_title = clean_title.replace(word, '').strip()
    
    # 첫 글자 대문자 처리 (영어인 경우)
    if clean_title and clean_title[0].isalpha():
        clean_title = clean_title[0].upper() + clean_title[1:]
    
    # 길이 제한
    if len(clean_title) > 60:
        clean_title = clean_title[:60] + "..."
    
    return clean_title

def extract_dates_from_text(text):
    """텍스트에서 날짜 추출"""
    dates = []
    
    for pattern in DATE_PATTERNS:
        matches = re.finditer(pattern, text)
        for match in matches:
            try:
                date_str = match.group(0)
                parsed_date = parse_korean_date(date_str)
                
                if parsed_date:
                    dates.append({
                        'original_text': date_str,
                        'parsed_date': parsed_date.isoformat(),
                        'confidence': 0.8
                    })
            except Exception as e:
                continue
    
    return dates

def extract_times_from_text(text):
    """텍스트에서 시간 추출"""  
    times = []
    
    for pattern in TIME_PATTERNS:
        matches = re.finditer(pattern, text)
        for match in matches:
            try:
                time_str = match.group(0)
                parsed_time = parse_korean_time(time_str)
                
                if parsed_time:
                    times.append({
                        'original_text': time_str,
                        'parsed_time': parsed_time,
                        'confidence': 0.8
                    })
            except Exception as e:
                continue
    
    return times

def parse_korean_date(date_str):
    """한국어 날짜 문자열을 datetime으로 변환"""
    try:
        # 상대적 날짜 처리
        today = datetime.now()
        
        if '오늘' in date_str:
            return today
        elif '내일' in date_str:
            return today + timedelta(days=1)
        elif '모레' in date_str:
            return today + timedelta(days=2)
        elif '다음주' in date_str:
            return today + timedelta(days=7)
        
        # 숫자 날짜 처리
        korean_date_match = re.search(r'(\d{4})?년?\s*(\d{1,2})월\s*(\d{1,2})일', date_str)
        if korean_date_match:
            year = korean_date_match.group(1) or today.year
            month = int(korean_date_match.group(2))
            day = int(korean_date_match.group(3))
            return datetime(int(year), month, day)
        
        # 다른 형식들도 시도
        return dateutil.parser.parse(date_str, fuzzy=True)
        
    except Exception as e:
        return None

def parse_korean_time(time_str):
    """한국어 시간 문자열 파싱"""
    try:
        # 오전/오후 처리
        if '오전' in time_str:
            hour_match = re.search(r'(\d{1,2})시', time_str)
            if hour_match:
                hour = int(hour_match.group(1))
                return f"{hour:02d}:00"
        
        elif '오후' in time_str:
            hour_match = re.search(r'(\d{1,2})시', time_str)
            if hour_match:
                hour = int(hour_match.group(1))
                if hour != 12:
                    hour += 12
                return f"{hour:02d}:00"
        
        # 24시간 형식
        time_match = re.search(r'(\d{1,2}):(\d{2})', time_str)
        if time_match:
            return f"{int(time_match.group(1)):02d}:{time_match.group(2)}"
        
        return None
        
    except Exception as e:
        return None

def deduplicate_todos_improved(todos):
    """개선된 중복 제거 - ID 기반"""
    seen_titles = set()
    unique_todos = []
    
    for todo in todos:
        title_key = todo['title'].lower()
        if title_key not in seen_titles:
            seen_titles.add(title_key)
            unique_todos.append(todo)
        else:
            print(f"[🗑️ 중복 제거] {todo['title']}")
    
    return unique_todos

def assign_priority(todos):
    """우선순위 자동 설정"""
    for todo in todos:
        # 키워드 기반 우선순위
        if todo['type'] == 'deadline':
            todo['priority'] = 'high'
        elif todo['type'] == 'meeting':
            todo['priority'] = 'high'
        elif '긴급' in todo['title'] or 'urgent' in todo['title'].lower():
            todo['priority'] = 'high'
        elif todo['type'] == 'event':
            todo['priority'] = 'medium'
        else:
            todo['priority'] = 'low'
        
        # 날짜 기반 우선순위 조정
        if todo['date']:
            try:
                todo_date = datetime.fromisoformat(todo['date'].replace('Z', '+00:00')).replace(tzinfo=None)
                days_until = (todo_date - datetime.now()).days
                
                if days_until <= 1:  # 오늘/내일
                    todo['priority'] = 'high'
                elif days_until <= 3:  # 3일 이내
                    if todo['priority'] == 'low':
                        todo['priority'] = 'medium'
            except:
                pass
    
    return todos
    
# ✅ 캐시 초기화 API 추가
@app.route('/api/clear-cache', methods=['POST'])
def clear_attachment_cache():
    """첨부파일 캐시 초기화"""
    global attachment_cache
    cache_count = len(attachment_cache)
    attachment_cache.clear()
    
    return jsonify({
        "success": True,
        "message": f"캐시 {cache_count}개 항목이 삭제되었습니다.",
        "cleared_items": cache_count
    })  

if __name__ == '__main__':
    print("=" * 60)
    print("🚀 YOLOv8 통합 메일 시스템 시작")
    print("=" * 60)
    
    # YOLO 모델 미리 로딩 (선택적)
    print("[🔄 YOLO 모델 사전 로딩 시도...]")
    load_yolo_model()
    
    print("=" * 60)
    app.run(debug=True, host="0.0.0.0", port=5001)
